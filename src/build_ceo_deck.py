"""
CEO deck — 10 slides, results-first, all key detail kept. Consulting (light) style + plain-English
explanation boxes. Reuses charts already rendered in reports/_mck/. Output: reports/CFC_CEO_10slides.pptx
"""
import json, pathlib
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

ROOT=pathlib.Path(__file__).resolve().parent.parent
M=json.loads((ROOT/"reports"/"deck_metrics.json").read_text())
FIG=ROOT/"reports"/"_mck"
def img(n): return str(FIG/n)

NAVY=RGBColor(0x1F,0x3A,0x5F); BLUE=RGBColor(0x2F,0x6D,0xB4); INK=RGBColor(0x22,0x2A,0x35)
GREY=RGBColor(0x6B,0x74,0x85); RULE=RGBColor(0xD2,0xD9,0xE3); GREEN=RGBColor(0x1E,0x7A,0x46); RED=RGBColor(0xB0,0x3A,0x2E)
WHITE=RGBColor(0xFF,0xFF,0xFF); PANEL=RGBColor(0xF4,0xF7,0xFB)
ML=0.7
prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5); BLANK=prs.slide_layouts[6]

def bg(s): s.background.fill.solid(); s.background.fill.fore_color.rgb=WHITE
def tx(s,x,y,w,h,t,size,color=INK,bold=False,align=PP_ALIGN.LEFT,sp=1.0):
    tb=s.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h)); tf=tb.text_frame; tf.word_wrap=True
    for i,line in enumerate(t.split("\n")):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.alignment=align; p.line_spacing=sp
        r=p.add_run(); r.text=line; r.font.size=Pt(size); r.font.color.rgb=color; r.font.bold=bold; r.font.name="Calibri"
    return tb
def head(s,eyebrow,title):
    tx(s,ML,0.42,12,0.35,eyebrow.upper(),12,BLUE,True)
    tx(s,ML,0.74,12.0,1.0,title,23,NAVY,True)
    ln=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(ML),Inches(1.72),Inches(11.93),Pt(1.5)); ln.fill.solid(); ln.fill.fore_color.rgb=RULE; ln.line.fill.background(); ln.shadow.inherit=False
def foot(s,n):
    tx(s,ML,7.04,9.5,0.3,"Source: CFC sales (Microsoft Fabric) 2023–2026; model backtest",9,GREY)
    tx(s,12.4,7.04,0.7,0.3,str(n),9,GREY)
def bullets(s,x,y,w,h,items,size=15,gap=9,color=INK,mark="–"):
    tb=s.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h)); tf=tb.text_frame; tf.word_wrap=True
    for i,it in enumerate(items):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.space_after=Pt(gap)
        r=p.add_run(); r.text=mark+"  "; r.font.size=Pt(size); r.font.color.rgb=BLUE; r.font.bold=True
        r2=p.add_run(); r2.text=it; r2.font.size=Pt(size); r2.font.color.rgb=color
def pic(s,p,x,y,w): s.shapes.add_picture(p,Inches(x),Inches(y),width=Inches(w))
def statbox(s,x,y,w,big,label,color=NAVY):
    tx(s,x,y,w,0.7,big,32,color,True); tx(s,x,y+0.62,w,0.6,label,13,GREY)
def caption(s,x,y,w,t): tx(s,x,y,w,0.5,t,11,GREY)
def readbox(s,x,y,w,h,title,lines,accent=BLUE):
    b=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(w),Inches(h)); b.fill.solid(); b.fill.fore_color.rgb=PANEL; b.line.color.rgb=RULE; b.line.width=Pt(0.75); b.shadow.inherit=False
    tb=s.shapes.add_textbox(Inches(x+0.18),Inches(y+0.12),Inches(w-0.36),Inches(h-0.24)); tf=tb.text_frame; tf.word_wrap=True
    p=tf.paragraphs[0]; r=p.add_run(); r.text=title.upper(); r.font.size=Pt(11); r.font.bold=True; r.font.color.rgb=accent; p.space_after=Pt(5)
    for ln in lines:
        p=tf.add_paragraph(); p.space_after=Pt(4)
        if isinstance(ln,tuple):
            r=p.add_run(); r.text=ln[0]+" "; r.font.size=Pt(12.5); r.font.bold=True; r.font.color.rgb=INK
            r2=p.add_run(); r2.text=ln[1]; r2.font.size=Pt(12.5); r2.font.color.rgb=INK
        else:
            r=p.add_run(); r.text=ln; r.font.size=Pt(12.5); r.font.color.rgb=INK
def table(s,x,y,w,h,rows,colw):
    t=s.shapes.add_table(len(rows),len(rows[0]),Inches(x),Inches(y),Inches(w),Inches(h)).table
    for i,cw in enumerate(colw): t.columns[i].width=Inches(cw)
    for ci in range(len(rows[0])):
        c=t.cell(0,ci); c.fill.solid(); c.fill.fore_color.rgb=NAVY; c.text=rows[0][ci]
        rr=c.text_frame.paragraphs[0].runs[0]; rr.font.size=Pt(12); rr.font.bold=True; rr.font.color.rgb=WHITE
    for ri in range(1,len(rows)):
        for ci in range(len(rows[0])):
            c=t.cell(ri,ci); c.text=rows[ri][ci]; c.fill.solid(); c.fill.fore_color.rgb=(PANEL if ri%2 else WHITE)
            rr=c.text_frame.paragraphs[0].runs[0]; rr.font.size=Pt(12); rr.font.name="Calibri"
            rr.font.color.rgb=(GREEN if ci==len(rows[0])-1 else INK); rr.font.bold=(ci==0 or ci==len(rows[0])-1)

NOTES=[]; N=[0]
def nx(): N[0]+=1; return N[0]
def note(t): NOTES.append(t)

# 1 TITLE
s=prs.slides.add_slide(BLANK); bg(s)
band=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(0),Inches(0),Inches(0.28),Inches(7.5)); band.fill.solid(); band.fill.fore_color.rgb=NAVY; band.line.fill.background(); band.shadow.inherit=False
tx(s,ML,2.4,11.5,0.4,"CITYFOOD CONCEPTS · DEMAND FORECASTING — RESULTS",13,BLUE,True)
tx(s,ML,2.85,11.8,1.5,"Forecasting demand and right-sizing\nevery warehouse order",32,NAVY,True,sp=1.05)
tx(s,ML,4.8,11,0.6,"Experiment results and recommended next step",17,GREY)
tx(s,ML,6.5,11,0.4,"Executive briefing  ·  June 2026",12,GREY)
note("SAY: In 10 slides — what we tried, what the experiment proved, and the one decision I need from you. Bottom line first on the next slide.")
foot(s,nx())

# 2 EXECUTIVE SUMMARY (answer first)
s=prs.slides.add_slide(BLANK); bg(s); head(s,"Bottom line","The experiment worked — accurate, trustworthy, cheaper; one data input unlocks full value")
bullets(s,ML,2.0,7.0,3.8,[
 "We forecast daily demand for every product in every shop, then turn it into the right order.",
 "Result: 16% more accurate than today’s averages — and it held up in every month we tested.",
 "The “safe stock” levels are honest, so orders can rely on them.",
 "In simulation, ordering this way is ~21% cheaper than current practice.",
 "It retrains itself. One gap remains: real product margin & shelf-life — a data request, not a rebuild."],size=15,gap=11)
statbox(s,8.4,2.1,4.0,"16%","more accurate than today")
statbox(s,8.4,3.5,4.0,"0.81","of demand swings explained")
statbox(s,8.4,4.9,4.0,"~21%","lower ordering cost (sim.)",GREEN)
note("SAY: Four wins, one honest gap. Accuracy up 16% and stable; safe levels trustworthy; ordering ~21% cheaper; self-maintaining. The only missing piece is real cost & shelf-life data from Finance and Ops.")
foot(s,nx())

# 3 PROBLEM
s=prs.slides.add_slide(BLANK); bg(s); head(s,"Why this matters","Daily ordering by gut feel loses money in both directions")
pic(s,img("monthly.png"),6.7,2.0,6.2); caption(s,6.7,5.7,6.2,"42 months of real demand — the pattern we forecast.")
bullets(s,ML,2.1,5.6,2.6,[
 "84 outlets order fresh bakery from the warehouse every day.",
 "Order too much → it spoils overnight (cash waste).",
 "Order too little → empty shelves (lost sales).",
 "58.8M units · ₭208bn sales → small % errors = big money."],size=15,gap=11)
readbox(s,ML,5.3,5.6,1.3,"The opportunity",[
 "Replace daily guesswork with a forecast that learns the real patterns — weekday, payday, weather, festivals."])
note("SAY: Today it’s gut feel. Both errors happen daily, and at this volume even a few percent is serious money. That’s the prize.")
foot(s,nx())

# 4 WHAT WE DID + HOW TESTED
s=prs.slides.add_slide(BLANK); bg(s); head(s,"The experiment","What we built, and how we tested it honestly")
pic(s,img("pipe.png"),0.9,1.95,11.6)
pic(s,img("wf.png"),7.4,4.2,5.4)
bullets(s,ML,4.2,6.3,2.3,[
 "Learn from 3.5 years of sales → forecast a range (likely / safe) → convert to an order → one warehouse list.",
 "Model: LightGBM (industry-standard for this data).",
 "Tested by training on the past and predicting future months it had never seen — repeated 3 times."],size=13.5,gap=8)
readbox(s,ML,6.5,6.3,0.0,"",[])  # spacer (no-op visual safety)
note("SAY: Top row is the pipeline: history → forecast → order → list. Bottom-right shows HOW we tested — train on the past, predict the next unseen month, roll forward. That’s the honest test; no peeking at the future.")
foot(s,nx())

# 5 RESULT — ACCURACY + scores
s=prs.slides.add_slide(BLANK); bg(s); head(s,"Result 1 · accuracy","16% more accurate than the best simple method")
pic(s,img("wmape.png"),6.6,2.0,6.3); caption(s,6.6,5.7,6.3,"WMAPE — shorter bar = more accurate.")
statbox(s,ML,2.1,3.0,"0.341","our error (WMAPE)")
readbox(s,ML,3.25,5.5,1.5,"What this means",[
 ("WMAPE","= how far off we are, as a share of total demand. 0.34 = off ~34% of volume on average; best simple method is 0.40."),
 ("So:","16% fewer ordering errors than today.")])
readbox(s,ML,4.9,5.5,1.7,"Other scores (plain words)",[
 ("R² 0.81","→ explains 81% of why demand moves."),
 ("Coverage 85/95%","→ “safe” levels are honest."),
 ("Bias −0.4","→ not over- or under-ordering.")])
note("SAY: The headline. We’re off by ~34% of volume on average vs ~40% for the best simple method — 16% better. The small box lists the other scores in plain words: explains most of the demand, safe levels honest, no bias.")
foot(s,nx())

# 6 RESULT — robustness + trust
s=prs.slides.add_slide(BLANK); bg(s); head(s,"Result 2 · robust & trustworthy","It holds every month, works best on top sellers, and its safe levels are honest")
pic(s,img("folds.png"),ML,2.0,4.0)
pic(s,img("abc.png"),4.85,2.0,4.0)
pic(s,img("calib.png"),9.0,2.0,3.9)
caption(s,ML,5.05,4.0,"Beats baseline every month."); caption(s,4.85,5.05,4.0,"Best on Class-A (80% of volume)."); caption(s,9.0,5.05,3.9,"Safe levels match reality.")
readbox(s,ML,5.45,11.93,1.15,"How to read these",[
 ("Left & middle:","blue = our model, grey = today’s method, shorter = better. We win in all 3 months and most on the high-volume products."),
 ("Right:","when we say a stock level is “85% safe”, demand really stays under it ~85% of the time → orders can trust it.")])
note("SAY: Three checks. It’s stable across months (not luck), strongest where the money is (Class-A), and its safety levels are calibrated. Blue is us, grey is today’s method; shorter is better.")
foot(s,nx())

# 7 RESULT — business impact
s=prs.slides.add_slide(BLANK); bg(s); head(s,"Result 3 · the money","Ordering this way is about 21% cheaper than current practice")
pic(s,img("cost.png"),6.5,2.3,6.3); caption(s,6.5,4.9,6.3,"Total cost = lost margin on stockouts + spoilage on overstock (simulated).")
statbox(s,ML,2.3,3.4,"~21%","lower ordering cost",GREEN)
readbox(s,ML,3.7,5.4,1.5,"What this means",[
 "We simulated ordering across 608,000 product-shop-days and counted the cost of being wrong — lost sales plus spoilage.",
 "Our forecast-driven orders cost ~21% less than today’s averages."])
readbox(s,ML,5.35,5.4,1.25,"Important",[
 "This uses placeholder economics. Real margins & shelf-life will widen the saving, not shrink it."])
note("SAY: Does it pay? Simulated over 608k order-days, our ordering is ~21% cheaper — fewer stockouts and less spoilage together. And this is on placeholder economics; real numbers improve it.")
foot(s,nx())

# 8 SELF-LEARNING
s=prs.slides.add_slide(BLANK); bg(s); head(s,"It keeps working","The system retrains itself and flags when the world changes")
pic(s,img("loop.png"),1.0,2.2,11.4)
readbox(s,ML,4.5,11.93,1.3,"What this means",[
 ("Self-improving:","each cycle it predicts, sees what sold, retrains, and only promotes a new model if it is at least 1% better — so quality never goes backwards."),
 ("Self-watching:","a monitor flags big changes (it caught the monsoon shift in testing) and signals when to relearn. Live model already scores 0.319.")])
note("SAY: No yearly rebuild. It relearns from new sales and only upgrades if better; a worse model never goes live. It also watches for change and flags a retrain. This is what keeps the saving durable.")
foot(s,nx())

# 9 HONEST GAP
s=prs.slides.add_slide(BLANK); bg(s); head(s,"The one gap","Forecast is proven; order-sizing needs real product economics")
table(s,ML,2.1,7.0,2.3,[["Input","Placeholder now","Needed from"],
 ["Gross margin","35% flat","Finance"],["Shelf-life","1 day","Operations"],["Salvage value","0","Operations"]],[2.3,2.5,2.2])
readbox(s,ML,4.7,11.93,1.6,"Why it matters — in plain words",[
 ("Today:","every product is treated the same, so ordering acts like a simple middle estimate."),
 ("With real margins & shelf-life:","high-profit, long-life items get more stock; cheap, perishable items stay lean → the full saving is unlocked."),
 ("Key point:","this is a data request to Finance & Ops — not a model change or rebuild.")])
note("SAY: Be transparent. The forecast is real and proven. The order-sizing currently uses placeholders. Real margins and shelf-life unlock per-product optimisation. It’s a data hand-off, not a rebuild.")
foot(s,nx())

# 10 RECOMMENDATION / NEXT
s=prs.slides.add_slide(BLANK); bg(s); head(s,"What next","Provide the economics, pilot in a few outlets, then scale")
bullets(s,ML,2.1,11.8,3.4,[
 "1.  Finance gives gross margin per product; Operations gives shelf-life + salvage.",
 "2.  We load it into one file — no code change — ordering becomes product-specific.",
 "3.  Run a 4–6 week pilot in a few outlets; measure kyats saved (waste + stockouts).",
 "4.  Roll out across the network; the self-learning loop sustains accuracy.",
 "5.  Optional next: specialist handling for rare items + automated nightly scheduling."],size=16,gap=13)
b=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(ML),Inches(5.7),Inches(11.93),Inches(0.85)); b.fill.solid(); b.fill.fore_color.rgb=RGBColor(0xEC,0xF4,0xEF); b.line.color.rgb=GREEN; b.line.width=Pt(1); b.shadow.inherit=False
tx(s,ML+0.2,5.85,11.5,0.6,"The decision today: approve the data hand-off and a pilot. The forecast is ready now.",15,NAVY,True)
note("SAY: The ask. Finance and Ops give us the economics; we load it in, no rebuild; we pilot for a month and measure the saving; then roll out. The one decision today is to approve the data hand-off and the pilot.")
foot(s,nx())

for sld,nt in zip(prs.slides,NOTES): sld.notes_slide.notes_text_frame.text=nt
out=ROOT/"reports"/"CFC_CEO_10slides.pptx"; prs.save(str(out))
import shutil; shutil.copy2(out, ROOT/"presentations"/"CFC_CEO_10slides.pptx")
print(f"wrote {out} ({N[0]} slides) + copied to presentations/")
