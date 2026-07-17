"""
Master evidence deck → native editable PowerPoint (.pptx).
Text is editable; charts rendered as PNG via matplotlib (dark theme, real numbers).
Output: reports/CFC_Model_Evidence_Scorecard.pptx
"""
import json, pathlib
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

ROOT=pathlib.Path(__file__).resolve().parent.parent
DD=json.loads((ROOT/"reports"/"deck_data.json").read_text())
M=json.loads((ROOT/"reports"/"deck_metrics.json").read_text())
FIG=ROOT/"reports"/"_pptx_figs"; FIG.mkdir(exist_ok=True)

# palette
BG=RGBColor(0x0E,0x14,0x20); CARD=RGBColor(0x17,0x22,0x34); INK=RGBColor(0xEE,0xF2,0xFA)
MUT=RGBColor(0x9B,0xAC,0xC8); ACC=RGBColor(0xFF,0x8A,0x4C); OK=RGBColor(0x3E,0xCF,0x8E)
BAD=RGBColor(0xE0,0x55,0x6B); LINE=RGBColor(0x26,0x34,0x4F)
HEXBG="#0e1420"; HEXACC="#ff8a4c"; HEXOK="#3ecf8e"; HEXMUT="#9bacc8"; HEXGREY="#37445f"; HEXBAD="#e0556b"; HEXGRID="#1e2942"

prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
BLANK=prs.slide_layouts[6]
SW,SH=prs.slide_width,prs.slide_height

def bg(s):
    s.background.fill.solid(); s.background.fill.fore_color.rgb=BG
def box(s,x,y,w,h,fill=CARD,line=LINE,lw=1.0):
    from pptx.enum.shapes import MSO_SHAPE
    sp=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(w),Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb=fill; sp.line.color.rgb=line; sp.line.width=Pt(lw)
    sp.shadow.inherit=False; return sp
def txt(s,x,y,w,h,runs,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,sp_after=6):
    tb=s.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h)); tf=tb.text_frame
    tf.word_wrap=True; tf.vertical_anchor=anchor
    if isinstance(runs,str): runs=[[(runs,18,INK,False)]]
    for i,para in enumerate(runs):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.alignment=align; p.space_after=Pt(sp_after)
        for (t,sz,c,b) in para:
            r=p.add_run(); r.text=t; r.font.size=Pt(sz); r.font.color.rgb=c; r.font.bold=b; r.font.name="Calibri"
    return tb
def header(s,h):
    txt(s,0.55,0.32,12.2,0.7,[[(h,26,INK,True)]])
def take(s,h_y,text,color=ACC):
    b=box(s,0.55,h_y,12.2,0.62,fill=RGBColor(0x1A,0x23,0x36),line=color,lw=1.5)
    tf=b.text_frame; tf.word_wrap=True; tf.margin_left=Inches(0.18); tf.vertical_anchor=MSO_ANCHOR.MIDDLE
    p=tf.paragraphs[0]; r=p.add_run(); r.text=text; r.font.size=Pt(15); r.font.bold=True
    r.font.color.rgb=RGBColor(0xFF,0xD9,0xBF) if color==ACC else RGBColor(0xFF,0xC9,0xD2)
def brand(s,n):
    txt(s,0.55,7.08,6,0.3,[[("CFC · CityFood Concepts",10,RGBColor(0x46,0x56,0x7A),False)]])
    txt(s,12.4,7.08,0.7,0.3,[[(str(n),10,RGBColor(0x46,0x56,0x7A),False)]])

# ---------- matplotlib charts ----------
def style(ax,fig):
    fig.patch.set_facecolor(HEXBG); ax.set_facecolor(HEXBG)
    ax.tick_params(colors=HEXMUT,labelsize=9)
    for sp in ax.spines.values(): sp.set_color(HEXGRID)
    ax.grid(color=HEXGRID,alpha=.5,linewidth=.6)
    if ax.get_legend(): [t.set_color(HEXMUT) for t in ax.get_legend().get_texts()]
def save(fig,name):
    p=FIG/f"{name}.png"; fig.savefig(p,dpi=150,bbox_inches="tight",facecolor=HEXBG); plt.close(fig); return str(p)

def c_scatter():
    fig,ax=plt.subplots(figsize=(5.3,3.4)); mx=DD["scatter_max"]
    xs=[d["x"] for d in DD["scatter"]]; ys=[d["y"] for d in DD["scatter"]]
    ax.scatter(xs,ys,s=5,c=HEXACC,alpha=.35); ax.plot([0,mx],[0,mx],"--",color=HEXOK)
    ax.set_xlim(0,mx); ax.set_ylim(0,mx); ax.set_xlabel("actual",color=HEXMUT); ax.set_ylabel("forecast",color=HEXMUT)
    style(ax,fig); return save(fig,"scatter")
def c_bar(labels,series,name,horiz=False,colors=None):
    fig,ax=plt.subplots(figsize=(5.3,3.3))
    import numpy as np; x=np.arange(len(labels)); w=.8/len(series)
    for i,(lab,data,col) in enumerate(series):
        if horiz: ax.barh(x,data,color=col,label=lab)
        else: ax.bar(x+i*w,data,w,color=(colors if colors else col),label=lab)
    if horiz:
        ax.set_yticks(x); ax.set_yticklabels(labels,fontsize=8); ax.invert_yaxis()
    else:
        ax.set_xticks(x+w*(len(series)-1)/2); ax.set_xticklabels(labels,fontsize=8,rotation=0)
    if len(series)>1: ax.legend(fontsize=8)
    style(ax,fig); return save(fig,name)
def c_line(labels,series,name):
    fig,ax=plt.subplots(figsize=(8.6,3.4))
    for lab,data,col,dash in series:
        ax.plot(range(len(labels)),data,color=col,label=lab,linewidth=1.6,linestyle=dash)
    step=max(1,len(labels)//8); ax.set_xticks(range(0,len(labels),step))
    ax.set_xticklabels([labels[i] for i in range(0,len(labels),step)],fontsize=8,rotation=0)
    ax.legend(fontsize=9); style(ax,fig); return save(fig,name)

IMG={
 "cScatter":c_scatter(),
 "cCalib":c_bar(["P50","P85","P95"],[("actual %",[M["cov50"],M["cov85"],M["cov95"]],HEXACC),("target %",[50,85,95],HEXGREY)],"calib"),
 "cWmape":c_bar(["LightGBM","7d-avg","28d-avg","naive","dow","wkday"],[("WMAPE",[0.341,0.405,0.411,0.453,0.468,0.537],None)],"wmape",
                colors=[HEXOK,HEXACC,HEXGREY,HEXGREY,HEXGREY,HEXGREY]),
 "cFolds":c_bar(["Apr-26","May-26","Jun-26"],[("LightGBM",[0.384,0.325,0.305],HEXOK),("baseline",[0.497,0.360,0.345],HEXGREY)],"folds"),
 "cAbc":c_bar(["A 75%","B 19%","C 5%"],[("LightGBM",[0.291,0.456,0.640],HEXOK),("baseline",[0.349,0.534,0.733],HEXGREY)],"abc"),
 "cResid":c_bar(M["resid_bins"],[("count",M["resid_counts"],HEXACC)],"resid"),
 "cHero":c_line(DD["hero"]["labels"],[("actual",DD["hero"]["actual"],"#eef2fa","-"),("P50",DD["hero"]["p50"],HEXACC,"-"),("P85",DD["hero"]["p85"],HEXOK,"--")],"hero"),
 "cImp":c_bar(["Product","Outlet","rmean28","yesterday","vol28","rmean7","month"],[("gain",[51107,30000,12101,6345,6100,5830,4203],HEXACC)],"imp",horiz=True),
 "cCost":c_bar(["baseline","P50","P85","P95"],[("cost ₭bn",[3.71,2.94,6.07,10.07],None)],"cost",colors=[HEXGREY,HEXOK,HEXACC,HEXBAD]),
}

# ---------- slide builders ----------
def s_title(t,sub,tag):
    s=prs.slides.add_slide(BLANK); bg(s)
    txt(s,1.0,2.4,11.3,1.6,[[(t,40,INK,True)]],align=PP_ALIGN.CENTER)
    txt(s,1.6,4.1,10.1,1.2,[[(sub,19,RGBColor(0xCD,0xD7,0xEA),False)]],align=PP_ALIGN.CENTER)
    txt(s,1.6,5.6,10.1,0.5,[[(tag,14,ACC,True)]],align=PP_ALIGN.CENTER)
    return s
def s_bullets(h,lead,bullets,warn=False,n=0):
    s=prs.slides.add_slide(BLANK); bg(s); header(s,h)
    txt(s,0.55,1.02,12.2,0.5,[[(lead,17,ACC,True)]])
    runs=[[(("⚠ " if warn else "▸ ")+b,19,(BAD if warn else INK) if False else INK,False)] for b in bullets]
    # bullet marker colored
    tb=s.shapes.add_textbox(Inches(0.7),Inches(1.7),Inches(11.9),Inches(5.0)); tf=tb.text_frame; tf.word_wrap=True
    for i,b in enumerate(bullets):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.space_after=Pt(10)
        r=p.add_run(); r.text=("⚠  " if warn else "▸  "); r.font.size=Pt(19); r.font.bold=True; r.font.color.rgb=(BAD if warn else ACC)
        r2=p.add_run(); r2.text=b; r2.font.size=Pt(19); r2.font.color.rgb=INK
    brand(s,n); return s
def s_table(h,take_txt,rows,headers=None,warn=False,n=0,colw=None):
    s=prs.slides.add_slide(BLANK); bg(s); header(s,h)
    y=1.05
    if take_txt: take(s,y,take_txt,BAD if warn else ACC); y=1.85
    nrows=len(rows)+(1 if headers else 0); ncols=len(rows[0])
    tb=s.shapes.add_table(nrows,ncols,Inches(0.55),Inches(y),Inches(12.2),Inches(min(5.2,0.55*nrows))).table
    if colw:
        for i,w in enumerate(colw): tb.columns[i].width=Inches(w)
    ri=0
    if headers:
        for ci,hh in enumerate(headers):
            c=tb.cell(0,ci); c.text=hh; c.fill.solid(); c.fill.fore_color.rgb=BG
            pr=c.text_frame.paragraphs[0]; pr.runs[0].font.size=Pt(12); pr.runs[0].font.bold=True; pr.runs[0].font.color.rgb=MUT
        ri=1
    for r in rows:
        for ci,val in enumerate(r):
            c=tb.cell(ri,ci); c.text=str(val); c.fill.solid(); c.fill.fore_color.rgb=CARD
            run=c.text_frame.paragraphs[0].runs[0]; run.font.size=Pt(14); run.font.name="Calibri"
            run.font.color.rgb=INK if ci==0 else RGBColor(0xCD,0xD7,0xEA)
            if ci==0: run.font.bold=True
        ri+=1
    brand(s,n); return s
def s_two(h,take_txt,left,right,warn=False,n=0):
    s=prs.slides.add_slide(BLANK); bg(s); header(s,h)
    take(s,1.05,take_txt,BAD if warn else ACC)
    lh,lp=left; rh,rp=right
    b1=box(s,0.55,2.0,5.9,4.4,line=BAD); b2=box(s,6.85,2.0,5.9,4.4,line=OK)
    txt(s,0.8,2.2,5.4,0.6,[[(lh,18,BAD,True)]]); txt(s,0.8,2.9,5.4,3.3,[[(lp,15,RGBColor(0xCD,0xD7,0xEA),False)]])
    txt(s,7.1,2.2,5.4,0.6,[[(rh,18,OK,True)]]); txt(s,7.1,2.9,5.4,3.3,[[(rp,15,RGBColor(0xCD,0xD7,0xEA),False)]])
    brand(s,n); return s
def s_chart(h,lead,img,wide=False,note=None,n=0):
    s=prs.slides.add_slide(BLANK); bg(s); header(s,h)
    txt(s,0.55,1.02,12.2,0.6,[[(lead,16,ACC,True)]])
    if wide: s.shapes.add_picture(img,Inches(1.4),Inches(1.9),width=Inches(10.5))
    else: s.shapes.add_picture(img,Inches(2.6),Inches(1.85),width=Inches(8.1))
    if note: txt(s,0.55,6.5,12.2,0.5,[[(note,13,MUT,False)]])
    brand(s,n); return s
def s_metric(h,take_txt,formula,value,read,good,img=None,note=None,n=0):
    s=prs.slides.add_slide(BLANK); bg(s); header(s,h); take(s,1.02,take_txt)
    colw=6.0 if img else 12.2
    fb=box(s,0.55,1.85,colw,0.75,fill=RGBColor(0x0A,0x0F,0x1A),line=ACC,lw=1.3)
    fb.text_frame.word_wrap=True; fb.text_frame.vertical_anchor=MSO_ANCHOR.MIDDLE
    p=fb.text_frame.paragraphs[0]; p.alignment=PP_ALIGN.CENTER; r=p.add_run(); r.text=formula; r.font.size=Pt(15); r.font.bold=True; r.font.color.rgb=ACC
    txt(s,0.55,2.75,colw,0.7,[[(value,30,INK,True)]])
    tb=s.shapes.add_textbox(Inches(0.6),Inches(3.55),Inches(colw-0.1),Inches(1.8)); tf=tb.text_frame; tf.word_wrap=True
    for i,b in enumerate(read):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.space_after=Pt(6)
        r=p.add_run(); r.text="▸  "; r.font.size=Pt(15); r.font.color.rgb=ACC; r.font.bold=True
        r2=p.add_run(); r2.text=b; r2.font.size=Pt(15); r2.font.color.rgb=INK
    gb=box(s,0.55,5.45,colw,0.8,fill=RGBColor(0x13,0x25,0x1C),line=OK)
    gb.text_frame.word_wrap=True; gb.text_frame.vertical_anchor=MSO_ANCHOR.MIDDLE; gb.text_frame.margin_left=Inches(0.15)
    p=gb.text_frame.paragraphs[0]; r=p.add_run(); r.text="Good looks like: "+good; r.font.size=Pt(13); r.font.color.rgb=RGBColor(0xBD,0xF0,0xD6)
    if img: s.shapes.add_picture(img,Inches(6.85),Inches(1.95),width=Inches(6.0))
    if note: txt(s,0.55,6.4,12.2,0.5,[[(note,12,MUT,False)]])
    brand(s,n); return s

# ---------- build deck (mirrors master HTML deck) ----------
n=0; nn=lambda: globals().__setitem__('n',n+1) or n+1
N=[0]
def k(): N[0]+=1; return N[0]

s_title("CFC Demand Forecasting — Model Evidence & Scorecard",
   "What we ran · what it produces · every score that applies · and how we know the predictions are right.",
   "Data Science Briefing · for CEO / CTO / DS · June 2026")
s_table("How to read this deck","Three layers on every slide — read the layer you need.",
   [["CEO","the orange takeaway line — the decision-relevant point"],
    ["CTO","the method & validation — is it sound, will it hold"],
    ["Data Scientist","the formula, value, target, interpretation"]],n=k(),colw=[3.0,9.2])
s_table("1 · The problem, defined precisely","We predict a number (units) per product per shop per day — a regression/forecasting task.",
   [["Target","net_units = gross − refunds − voids (true demand)"],["Grain","one prediction per (outlet × product × day)"],
    ["Horizon","next-day (extendable)"],["Type","regression + quantile forecasting (NOT classification)"],
    ["Universe","Class A+B finished goods · 402 products × 84 outlets"]],n=k(),colw=[3.0,9.2])
s_table("2 · What we ran — model specification","A gradient-boosted tree model (LightGBM) — proven workhorse for tabular demand.",
   [["Algorithm","LightGBM, quantile objective (P50 / P85 / P95)"],["Trees / leaves","600 estimators · 255 leaves · lr 0.05"],
    ["Regularisation","min_child 100 · subsample 0.8 · colsample 0.8"],["Features","37 (18 categorical incl Product & Outlet IDs)"],
    ["Routing","smooth/erratic→LightGBM; intermittent/lumpy→Croston (backlog)"]],n=k(),colw=[3.0,9.2])
s_table("3 · Data & split","Trained on the past, tested on the future it never saw — like production.",
   [["Training rows","7,721,315 (before Apr-2026)"],["Test rows","608,261 (Apr–Jun 2026, held out)"],
    ["Validation","rolling-origin walk-forward, 3 monthly folds"],["Leak control","history features shifted ≥1 day; no future info"],
    ["Champion (live)","retrained on 7.9M, 60-day holdout"]],n=k(),colw=[3.0,9.2])
s_bullets("4 · Validation method — why it proves real-world skill","Walk-forward backtest: the honest way to test a forecaster.",
   ["Train up to month M, predict M+1, score, roll forward.","Model NEVER sees the test period in training — no leakage.",
    "Repeated over 3 months → tests stability, not luck.","Mirrors exactly how it runs in production each night.",
    "A single random split would overstate accuracy for time series — we avoid it."],n=k())
s_two("5 · Which metrics apply — and which DON'T","AUC / ROC / precision / recall are classification metrics. They do NOT apply to this regression task.",
   ("NOT used (classification only)","AUC, ROC curve, precision, recall, F1, confusion matrix, log-loss.\n\nThese need yes/no labels. We predict a quantity, not a class. Reporting them here would be wrong."),
   ("USED (regression + quantile)","WMAPE, MAE, RMSE, Bias, R² — for the point forecast.\nPinball loss + Coverage — for the P85/P95 safety levels.\nSkill vs baseline — proves it beats simple methods."),warn=True,n=k())

s_metric("6 · WMAPE — headline accuracy","On average our forecast is off by 34% of volume — 16% better than today's method.",
   "WMAPE = Σ|actual − forecast| / Σ|actual|",str(M["wmape"]),
   ["Total miss in units ÷ total real demand. Volume-weighted.","Stable where plain MAPE blows up (near-zero days)."],
   "Retail demand: <0.40 workable, <0.35 good. We are 0.341.",n=k())
s_metric("7 · MAE — average miss in units","Typical daily miss ≈ 2.3 units per product-shop line.",
   "MAE = mean( |actual − forecast| )",f"{M['mae']} units",
   ["Plain-English error size, same units as demand.","Robust to outliers (unlike RMSE)."],
   "Lower is better; judge vs the item's typical daily volume.",n=k())
s_metric("8 · RMSE — penalises big misses","Large errors are rare; RMSE near MAE means few wild misses.",
   "RMSE = sqrt( mean( (actual − forecast)² ) )",f"{M['rmse']} units",
   ["Squares errors → punishes big mistakes harder than MAE.",f"RMSE {M['rmse']} vs MAE {M['mae']}: spread is controlled."],
   "Want RMSE not much larger than MAE → no catastrophes.",n=k())
s_metric("9 · Bias — over or under forecasting?","Almost unbiased — tiny tendency to under-forecast (0.41 u).",
   "Bias = mean( forecast − actual )",f"{M['bias']} units",
   ["Near 0 = balanced. Negative = slight under-forecast.","Persistent bias would systematically over/under-stock."],
   "Want ≈0. −0.41 on ~6-unit orders is negligible.",n=k())
s_metric("10 · R² — variance explained","The model explains 81% of the variation in demand.",
   "R² = 1 − SS(residual) / SS(total)",str(M["r2"]),
   ["1.0 perfect, 0 = no better than the average.","0.811 → captures most day-to-day demand swing."],
   ">0.7 strong for noisy daily retail. We are 0.81.",img=IMG["cScatter"],
   note="Scatter: each dot = one outlet-product-day; tight along green diagonal = accurate.",n=k())
s_metric("11 · Pinball loss — quantile quality","P50/P85/P95 scored with the metric built for quantiles.",
   "Pinball(q) = mean( max( q·e , (q−1)·e ) )",f"P50 {M['pin50']} · P85 {M['pin85']} · P95 {M['pin95']}",
   ["Correct loss for quantile forecasts (asymmetric).","Lower = sharper AND correctly-placed intervals."],
   "Used to compare quantile models; lower better.",n=k())
s_metric("12 · Coverage / Calibration — honest safety levels","When we say '85% safe', demand stays under it 85.4% of the time.",
   "Coverage(q) = share of actuals ≤ forecast at q",f"P85 {M['cov85']}% · P95 {M['cov95']}% (P50 {M['cov50']}%)",
   ["P85/P95 on target → intervals trustworthy.","P50 56.7% reflects small under-bias (orders round up safely)."],
   "Coverage ≈ stated quantile = calibrated → safe buffers.",img=IMG["cCalib"],n=k())

s_chart("13 · Proof 1 — beats the baseline (skill)","Only 'good' if it beats simple methods on the SAME test. It does, +16%.",IMG["cWmape"],
   note="LightGBM 0.341 vs best simple method (7-day avg) 0.405. Lower = better.",n=k())
s_chart("14 · Proof 2 — stable across folds (not luck/overfit)","Beats the floor in all 3 test months — consistent.",IMG["cFolds"],n=k())
s_chart("15 · Proof 3 — accurate where it matters (ABC)","Best on Class-A (80% of volume): WMAPE 0.291.",IMG["cAbc"],n=k())
s_chart("16 · Proof 4 — residuals centred & tight","Errors cluster at zero with thin tails → no systematic mistake.",IMG["cResid"],
   note="Distribution of (actual − forecast). Symmetric, peaked at 0 = healthy.",n=k())
s_chart("17 · Proof 5 — tracks reality (hero product)",f"{DD['hero']['name']}: forecast (orange) follows actual (white).",IMG["cHero"],wide=True,n=k())
s_chart("18 · Proof 6 — leans on sensible drivers","Top signals = product, outlet, recent rolling demand. No leakage artefacts.",IMG["cImp"],n=k())

s_bullets("19 · \"How do I KNOW the prediction is right?\"","Six independent checks — the answer you can give with confidence.",
   ["Skill: beats best simple method by +16% on unseen data.","Stability: wins all 3 separate test months.",
    "Calibration: 'X% safe' is right X% of the time (85.4 / 94.5%).","Low bias: −0.41 u → not systematically over/under.",
    "Variance explained: R² 0.81.","Business test: cuts simulated ordering cost ~21%."],n=k())
s_bullets("20 · The outcome — what the model produces","Predictions become decisions.",
   ["Demand range (P50/P85/P95) per product per outlet per day.","Order quantity via newsvendor (stockout vs spoilage).",
    "One daily warehouse picklist (~230 products, ~34,000 units).","A service-level dial for management policy.","Refreshed automatically every night."],n=k())
s_chart("21 · Business validation — cost simulation","Lower ordering cost than current practice. Model P50 ≈ 21% cheaper.",IMG["cCost"],
   note="Simulated over 608k order-days. Demo economics; real margin+shelf-life will improve it.",n=k())
s_table("22 · Self-learning & drift control","Stays accurate over time, on its own, with guardrails.",
   [["Champion / challenger","new model promoted only if ≥1% better"],["Champion WMAPE (holdout)","0.319 — beats stretch target ≤0.321"],
    ["Drift monitor","PSI on inputs + accuracy creep → retrain signal"],["Caught in test","monsoon weather shift flagged; accuracy held"]],n=k(),colw=[3.4,8.8])
s_table("23 · Assumptions & limitations (honest)","Forecast is real and proven. Order-sizing economics are placeholders until Finance/Ops data arrives.",
   [["Assumed","margin 35% flat · shelf-life 1 day · salvage 0"],["Effect","flat critical ratio → ordering ≈ P50 (no per-product edge yet)"],
    ["Not modelled yet","intermittent tail (Croston), promo uplift detail"],["Data gap, not model gap","real economics = a data request, no rebuild"]],warn=True,n=k(),colw=[3.4,8.8])
s_bullets("24 · What would prove us WRONG (falsifiability)","A credible model states how it could fail — and we monitor for it.",
   ["Accuracy on new months drifts above ~0.40 WMAPE → retrain trigger.","Coverage drifts off target → recalibrate.",
    "Bias grows persistently → investigate features.","Drift monitor (PSI) flags unseen input shift → relearn.","All watched automatically; none silent."],n=k())
s_table("25 · The scorecard — every metric at a glance",None,
   [["WMAPE",str(M["wmape"]),"<0.40 / tgt 0.321","good (+16%)"],["MAE",f"{M['mae']} u","vs ~6u order","small"],
    ["RMSE",f"{M['rmse']} u","near MAE","no outliers"],["Bias",f"{M['bias']} u","≈0","near-neutral"],
    ["R²",str(M["r2"]),">0.70","strong"],["Pinball P85",str(M["pin85"]),"lower better","sharp"],
    ["Coverage P85",f"{M['cov85']}%","≈85%","calibrated"],["Coverage P95",f"{M['cov95']}%","≈95%","calibrated"],
    ["Champion holdout","0.319","≤0.321","beats target"],["Cost vs baseline","−21%","<0","saves money"]],
   headers=["Metric","Value","Target / ref","Verdict"],n=k(),colw=[3.2,2.3,3.7,3.0])
s_bullets("26 · Next step — the one unlock","Real economics turns a proven forecast into optimal per-product ordering.",
   ["Finance: gross margin per product. Ops: shelf-life + salvage.","Load into data/product_econ.csv — no code change.",
    "Per-product critical ratios spread → smart safety stock.","Add Croston for rare items; schedule nightly/weekly jobs.",
    "Pilot in a few outlets, measure ₭ saved, then roll out."],n=k())
s_table("27 · Appendix — metric glossary (one line each)",None,
   [["WMAPE","weighted avg % error; headline accuracy"],["MAE","mean absolute error, in units"],
    ["RMSE","root mean squared error; penalises big misses"],["Bias","mean(forecast−actual); over/under tendency"],
    ["R²","fraction of demand variance explained"],["Pinball loss","correct loss for quantile (P85/P95) forecasts"],
    ["Coverage","how often actual ≤ the quantile (calibration)"],["PSI","population stability index; input-drift detector"],
    ["AUC/ROC","classification-only — NOT used (this is regression)"]],n=k(),colw=[3.0,9.2])

out=ROOT/"reports"/"CFC_Model_Evidence_Scorecard.pptx"
prs.save(str(out))
print(f"wrote {out} — {len(prs.slides.__iter__.__self__._sldIdLst)} slides")
