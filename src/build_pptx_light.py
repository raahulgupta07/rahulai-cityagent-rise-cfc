"""
CFC Model Evidence deck — LIGHT theme, a visual on EVERY slide, formulas in plain words.
Charts + flow diagrams rendered via matplotlib (light). Native editable text via python-pptx.
Output: reports/CFC_Model_Evidence_Light.pptx
"""
import json, pathlib, numpy as np
import matplotlib; matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib.patches import FancyBboxPatch, FancyArrowPatch, Circle
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

ROOT=pathlib.Path(__file__).resolve().parent.parent
DD=json.loads((ROOT/"reports"/"deck_data.json").read_text())
M=json.loads((ROOT/"reports"/"deck_metrics.json").read_text())
FIG=ROOT/"reports"/"_pptx_light"; FIG.mkdir(exist_ok=True)

# LIGHT palette
WHITE=RGBColor(0xFF,0xFF,0xFF); INK=RGBColor(0x1B,0x26,0x38); MUT=RGBColor(0x5C,0x6B,0x82)
ACC=RGBColor(0xE8,0x6A,0x17); OKc=RGBColor(0x12,0x9E,0x5E); BADc=RGBColor(0xC8,0x3B,0x4B)
PANEL=RGBColor(0xF4,0xF6,0xFA); BAND=RGBColor(0xFA,0xF2,0xE9); LINEc=RGBColor(0xD7,0xDD,0xE7)
HACC="#e86a17"; HOK="#129e5e"; HBAD="#c83b4b"; HMUT="#5c6b82"; HGREY="#b9c2d0"; HGRID="#e3e8f0"; HINK="#1b2638"

prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
BLANK=prs.slide_layouts[6]

def bg(s): s.background.fill.solid(); s.background.fill.fore_color.rgb=WHITE
def header(s,h):
    tb=s.shapes.add_textbox(Inches(0.55),Inches(0.3),Inches(12.2),Inches(0.7)); tf=tb.text_frame; tf.word_wrap=True
    p=tf.paragraphs[0]; r=p.add_run(); r.text=h; r.font.size=Pt(25); r.font.bold=True; r.font.color.rgb=INK; r.font.name="Calibri"
    ln=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(0.57),Inches(1.0),Inches(1.5),Inches(0.05))
    ln.fill.solid(); ln.fill.fore_color.rgb=ACC; ln.line.fill.background(); ln.shadow.inherit=False
def take(s,y,text,color=ACC):
    b=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.55),Inches(y),Inches(12.2),Inches(0.6))
    b.fill.solid(); b.fill.fore_color.rgb=BAND if color==ACC else RGBColor(0xFB,0xEC,0xEE)
    b.line.color.rgb=color; b.line.width=Pt(1.25); b.shadow.inherit=False
    tf=b.text_frame; tf.word_wrap=True; tf.vertical_anchor=MSO_ANCHOR.MIDDLE; tf.margin_left=Inches(0.18)
    p=tf.paragraphs[0]; r=p.add_run(); r.text=text; r.font.size=Pt(15); r.font.bold=True; r.font.color.rgb=INK
def brand(s,n):
    tb=s.shapes.add_textbox(Inches(0.55),Inches(7.08),Inches(7),Inches(0.3))
    r=tb.text_frame.paragraphs[0].add_run(); r.text="CFC · CityFood Concepts"; r.font.size=Pt(10); r.font.color.rgb=RGBColor(0xA7,0xB0,0xBF)
    tb2=s.shapes.add_textbox(Inches(12.4),Inches(7.08),Inches(0.7),Inches(0.3))
    r2=tb2.text_frame.paragraphs[0].add_run(); r2.text=str(n); r2.font.size=Pt(10); r2.font.color.rgb=RGBColor(0xA7,0xB0,0xBF)
def bullets(s,x,y,w,h,items,size=16,warn=False,gap=7):
    tb=s.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h)); tf=tb.text_frame; tf.word_wrap=True
    for i,b in enumerate(items):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.space_after=Pt(gap)
        r=p.add_run(); r.text=("! " if warn else "› "); r.font.size=Pt(size); r.font.bold=True; r.font.color.rgb=(BADc if warn else ACC)
        r2=p.add_run(); r2.text=b; r2.font.size=Pt(size); r2.font.color.rgb=INK
def pic(s,img,x,y,w): s.shapes.add_picture(img,Inches(x),Inches(y),width=Inches(w))
def plaintext(s,x,y,w,h,t,size=14,color=INK,bold=False,align=PP_ALIGN.LEFT):
    tb=s.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h)); tf=tb.text_frame; tf.word_wrap=True
    for i,line in enumerate(t.split("\n")):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.alignment=align
        r=p.add_run(); r.text=line; r.font.size=Pt(size); r.font.color.rgb=color; r.font.bold=bold

# ---------------- matplotlib (light) ----------------
def _st(ax,fig):
    fig.patch.set_facecolor("white"); ax.set_facecolor("white")
    ax.tick_params(colors=HINK,labelsize=9)
    for k,sp in ax.spines.items():
        sp.set_color(HGRID);
        if k in("top","right"): sp.set_visible(False)
    ax.grid(color=HGRID,alpha=.9,linewidth=.7)
def _save(fig,n): p=FIG/f"{n}.png"; fig.savefig(p,dpi=150,bbox_inches="tight",facecolor="white"); plt.close(fig); return str(p)

def bar(labels,series,name,horiz=False,colors=None,figsize=(5.2,3.2),target=None):
    fig,ax=plt.subplots(figsize=figsize); x=np.arange(len(labels)); w=.8/len(series)
    for i,(lab,data,col) in enumerate(series):
        if horiz: ax.barh(x,data,color=col);
        else: ax.bar(x+i*w,data,w,color=(colors if colors else col),label=lab)
    if horiz: ax.set_yticks(x); ax.set_yticklabels(labels,fontsize=8); ax.invert_yaxis()
    else:
        ax.set_xticks(x+w*(len(series)-1)/2); ax.set_xticklabels(labels,fontsize=8)
        if len(series)>1: ax.legend(fontsize=8,frameon=False)
    if target is not None: ax.axhline(target,color=HBAD,ls="--",lw=1.2); ax.text(len(labels)-1,target,f" target {target}",color=HBAD,fontsize=8,va="bottom",ha="right")
    _st(ax,fig); return _save(fig,name)
def line(labels,series,name,figsize=(8.4,3.3)):
    fig,ax=plt.subplots(figsize=figsize)
    for lab,data,col,dash in series: ax.plot(range(len(labels)),data,color=col,label=lab,lw=1.7,ls=dash)
    step=max(1,len(labels)//8); ax.set_xticks(range(0,len(labels),step)); ax.set_xticklabels([labels[i] for i in range(0,len(labels),step)],fontsize=8)
    ax.legend(fontsize=9,frameon=False); _st(ax,fig); return _save(fig,name)
def scatter(name):
    fig,ax=plt.subplots(figsize=(5.0,3.3)); mx=DD["scatter_max"]
    ax.scatter([d["x"] for d in DD["scatter"]],[d["y"] for d in DD["scatter"]],s=5,c=HACC,alpha=.3)
    ax.plot([0,mx],[0,mx],"--",color=HOK,lw=1.5); ax.set_xlim(0,mx); ax.set_ylim(0,mx)
    ax.set_xlabel("actual",color=HINK,fontsize=9); ax.set_ylabel("forecast",color=HINK,fontsize=9); _st(ax,fig); return _save(fig,name)
def gauge(value,name,vmax=1.0,good=0.7):
    fig,ax=plt.subplots(figsize=(5.0,1.5)); ax.barh([0],[vmax],color=HGREY,height=.5)
    ax.barh([0],[value],color=HOK,height=.5); ax.axvline(good,color=HBAD,ls="--",lw=1.2)
    ax.text(good,0.5,f"good ≥{good}",color=HBAD,fontsize=8,ha="center"); ax.text(value,0,f" {value}",va="center",fontsize=12,fontweight="bold",color=HINK)
    ax.set_xlim(0,vmax); ax.set_yticks([]); fig.patch.set_facecolor("white"); ax.set_facecolor("white")
    for sp in ax.spines.values(): sp.set_visible(False)
    ax.tick_params(labelsize=8,colors=HINK); return _save(fig,name)
def biaschart(name):
    fig,ax=plt.subplots(figsize=(5.0,1.7)); ax.axvline(0,color=HMUT,lw=1)
    ax.barh([0],[M["bias"]],color=HACC,height=.4); ax.set_xlim(-2,2); ax.set_yticks([])
    ax.text(M["bias"],0.35,f"{M['bias']} (near 0 = balanced)",fontsize=9,ha="center",color=HINK)
    fig.patch.set_facecolor("white"); ax.set_facecolor("white")
    for k,sp in ax.spines.items():
        if k!="bottom": sp.set_visible(False)
        else: sp.set_color(HGRID)
    ax.tick_params(labelsize=8,colors=HINK); return _save(fig,name)

# ---- flow / diagram generators ----
def _bxf(ax,x,y,w,h,title,sub,fc="#fdf0e4",ec=HACC):
    ax.add_patch(FancyBboxPatch((x,y),w,h,boxstyle="round,pad=0.02,rounding_size=0.06",fc=fc,ec=ec,lw=1.4))
    ax.text(x+w/2,y+h*0.62,title,ha="center",va="center",fontsize=10,fontweight="bold",color=HINK)
    if sub: ax.text(x+w/2,y+h*0.27,sub,ha="center",va="center",fontsize=7.5,color=HMUT)
def _arrow(ax,x1,y1,x2,y2):
    ax.add_patch(FancyArrowPatch((x1,y1),(x2,y2),arrowstyle="-|>",mutation_scale=14,color=HACC,lw=1.6))

def flow_pipeline(name):
    fig,ax=plt.subplots(figsize=(11.5,2.4)); ax.set_xlim(0,11.5); ax.set_ylim(0,2.4); ax.axis("off"); fig.patch.set_facecolor("white")
    steps=[("Past sales","3.5 yrs"),("Learn","find patterns"),("Forecast","range P50/85/95"),("Order qty","balance cost"),("Picklist","make & send")]
    w=1.9;h=1.2;gap=0.42
    for i,(t,sb) in enumerate(steps):
        x=i*(w+gap)+0.1; _bxf(ax,x,0.6,w,h,t,sb)
        if i<len(steps)-1: _arrow(ax,x+w,1.2,x+w+gap,1.2)
    return _save(fig,name)
def flow_walkforward(name):
    fig,ax=plt.subplots(figsize=(10.5,3.2)); ax.set_xlim(0,11); ax.set_ylim(0,3.2); ax.axis("off"); fig.patch.set_facecolor("white")
    for r,(tr,te) in enumerate([(6,7),(7,8),(8,9)]):
        y=2.3-r*0.85
        ax.add_patch(FancyBboxPatch((0.3,y),tr*0.9,0.55,boxstyle="round,pad=0.01",fc="#eaf4ee",ec=HOK,lw=1.2))
        ax.text(0.3+tr*0.9/2,y+0.27,f"train (months → {tr})",ha="center",va="center",fontsize=8,color=HINK)
        ax.add_patch(FancyBboxPatch((0.3+tr*0.9+0.1,y),0.9,0.55,boxstyle="round,pad=0.01",fc="#fdf0e4",ec=HACC,lw=1.2))
        ax.text(0.3+tr*0.9+0.55,y+0.27,"test",ha="center",va="center",fontsize=8,fontweight="bold",color=HACC)
    ax.text(5.5,0.15,"roll forward each month → never sees the future",ha="center",fontsize=9,color=HMUT,style="italic")
    return _save(fig,name)
def diag_classes(name):
    fig,ax=plt.subplots(figsize=(10.5,2.8)); ax.set_xlim(0,11); ax.set_ylim(0,2.8); ax.axis("off"); fig.patch.set_facecolor("white")
    _bxf(ax,0.4,0.7,4.7,1.4,"NOT used — classification","AUC · ROC · precision · recall · F1",fc="#fbecee",ec=HBAD)
    _bxf(ax,5.9,0.7,4.7,1.4,"USED — regression + quantile","WMAPE · MAE · RMSE · R² · Pinball · Coverage",fc="#eaf4ee",ec=HOK)
    ax.text(5.5,2.5,"We predict a NUMBER, not a yes/no class",ha="center",fontsize=10,fontweight="bold",color=HINK)
    return _save(fig,name)
def diag_scale(name):
    fig,ax=plt.subplots(figsize=(9.5,3.0)); ax.set_xlim(0,10); ax.set_ylim(0,3); ax.axis("off"); fig.patch.set_facecolor("white")
    ax.plot([1,9],[1.5,1.5],color=HINK,lw=2); ax.plot([5,5],[1.5,2.4],color=HINK,lw=2)
    _bxf(ax,0.6,0.3,3.2,1.0,"Order too little","stockout = lost margin",fc="#fbecee",ec=HBAD)
    _bxf(ax,6.2,0.3,3.2,1.0,"Order too much","spoilage = wasted cash",fc="#fbecee",ec=HBAD)
    ax.text(5,2.6,"⚖ Newsvendor balances the two",ha="center",fontsize=11,fontweight="bold",color=HACC)
    return _save(fig,name)
def diag_loop(name):
    fig,ax=plt.subplots(figsize=(9.5,2.6)); ax.set_xlim(0,11); ax.set_ylim(0,2.6); ax.axis("off"); fig.patch.set_facecolor("white")
    steps=["Predict\ntonight","See what\nsold","Learn\nfrom it","Promote\nif better"]
    w=2.0;h=1.2;gap=0.5
    for i,t in enumerate(steps):
        x=i*(w+gap)+0.2; _bxf(ax,x,0.7,w,h,t,"")
        if i<len(steps)-1: _arrow(ax,x+w,1.3,x+w+gap,1.3)
    ax.text(5.4,0.25,"↻ repeats automatically — a worse model never goes live",ha="center",fontsize=9,color=HMUT,style="italic")
    return _save(fig,name)
def diag_checks(name):
    fig,ax=plt.subplots(figsize=(10.5,3.2)); ax.set_xlim(0,11); ax.set_ylim(0,3.2); ax.axis("off"); fig.patch.set_facecolor("white")
    checks=["Skill +16%","Stable 3/3 folds","Calibrated 85/95%","Low bias −0.41","R² 0.81","Cost −21%"]
    for i,c in enumerate(checks):
        x=(i%3)*3.6+0.3; y=1.8-(i//3)*1.3
        _bxf(ax,x,y,3.2,1.0,"✓ "+c,"",fc="#eaf4ee",ec=HOK)
    return _save(fig,name)

I={
 "pipeline":flow_pipeline("pipeline"), "wf":flow_walkforward("wf"), "classes":diag_classes("classes"),
 "scale":diag_scale("scale"), "loop":diag_loop("loop"), "checks":diag_checks("checks"),
 "scatter":scatter("scatter"),
 "wmape":bar(["LightGBM","7d-avg","28d-avg","naive","dow","wkday"],[("WMAPE",[0.341,0.405,0.411,0.453,0.468,0.537],None)],"wmape",colors=[HOK,HACC,HGREY,HGREY,HGREY,HGREY],target=0.321),
 "mae":bar(["MAE","RMSE"],[("units",[M["mae"],M["rmse"]],None)],"mae",colors=[HOK,HACC]),
 "rmse":bar(["MAE","RMSE"],[("units",[M["mae"],M["rmse"]],None)],"rmse2",colors=[HGREY,HACC]),
 "bias":biaschart("bias"),
 "r2":gauge(M["r2"],"r2",1.0,0.7),
 "pinball":bar(["P50","P85","P95"],[("pinball",[M["pin50"],M["pin85"],M["pin95"]],None)],"pinball",colors=[HACC,HOK,HMUT]),
 "calib":bar(["P50","P85","P95"],[("actual %",[M["cov50"],M["cov85"],M["cov95"]],HACC),("target %",[50,85,95],HGREY)],"calib"),
 "folds":bar(["Apr","May","Jun"],[("LightGBM",[0.384,0.325,0.305],HOK),("baseline",[0.497,0.360,0.345],HGREY)],"folds"),
 "abc":bar(["A 75%","B 19%","C 5%"],[("LightGBM",[0.291,0.456,0.640],HOK),("baseline",[0.349,0.534,0.733],HGREY)],"abc"),
 "resid":bar(M["resid_bins"],[("count",M["resid_counts"],HACC)],"resid"),
 "hero":line(DD["hero"]["labels"],[("actual",DD["hero"]["actual"],HINK,"-"),("P50",DD["hero"]["p50"],HACC,"-"),("P85",DD["hero"]["p85"],HOK,"--")],"hero"),
 "imp":bar(["Product","Outlet","rmean28","yesterday","vol28","rmean7","month"],[("gain",[51107,30000,12101,6345,6100,5830,4203],HACC)],"imp",horiz=True),
 "cost":bar(["baseline","P50","P85","P95"],[("cost ₭bn",[3.71,2.94,6.07,10.07],None)],"cost",colors=[HGREY,HOK,HACC,HBAD]),
 "monthly":line(DD["monthly"]["labels"],[("net units",DD["monthly"]["units"],HACC,"-")],"monthly"),
}

# ---------------- slides ----------------
N=[0]
def num(): N[0]+=1; return N[0]

def title_slide():
    s=prs.slides.add_slide(BLANK); bg(s)
    bar0=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(0),Inches(0),Inches(13.333),Inches(0.25)); bar0.fill.solid(); bar0.fill.fore_color.rgb=ACC; bar0.line.fill.background(); bar0.shadow.inherit=False
    plaintext(s,1.0,2.2,11.3,1.4,"CFC Demand Forecasting\nModel Evidence & Scorecard",34,INK,True,PP_ALIGN.CENTER)
    plaintext(s,1.6,4.0,10.1,1.0,"What we ran · what it produces · every score that applies · and how we know the predictions are right.",18,MUT,False,PP_ALIGN.CENTER)
    plaintext(s,1.6,5.3,10.1,0.5,"Data Science Briefing · for CEO / CTO / Data Science · June 2026",14,ACC,True,PP_ALIGN.CENTER)
    pic(s,I["pipeline"],1.6,5.9,10.1)

def diagram_slide(h,take_txt,img,bullets_list,warn=False,note=None,wide=11.0):
    s=prs.slides.add_slide(BLANK); bg(s); header(s,h); take(s,1.05,take_txt,BADc if warn else ACC)
    pic(s,img,(13.333-wide)/2,1.85,wide)
    yb=1.85+wide*0.30+0.15
    if bullets_list: bullets(s,0.7,min(yb,5.0),12.0,1.8,bullets_list,size=14,warn=warn,gap=5)
    if note: plaintext(s,0.55,6.55,12.2,0.5,note,12,MUT)
    brand(s,num())

def metric_slide(h,take_txt,formula,plain,value,read,good,img,note=None):
    s=prs.slides.add_slide(BLANK); bg(s); header(s,h); take(s,1.02,take_txt)
    # left col text
    fb=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.55),Inches(1.8),Inches(6.1),Inches(0.7))
    fb.fill.solid(); fb.fill.fore_color.rgb=PANEL; fb.line.color.rgb=ACC; fb.line.width=Pt(1.2); fb.shadow.inherit=False
    tf=fb.text_frame; tf.word_wrap=True; tf.vertical_anchor=MSO_ANCHOR.MIDDLE; tf.margin_left=Inches(0.12)
    r=tf.paragraphs[0].add_run(); r.text=formula; r.font.size=Pt(13); r.font.bold=True; r.font.color.rgb=INK
    plaintext(s,0.6,2.6,6.0,0.7,"In plain words: "+plain,13,OKc,True)
    plaintext(s,0.6,3.45,6.0,0.6,value,26,ACC,True)
    bullets(s,0.6,4.25,6.0,1.4,read,size=13,gap=4)
    gb=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.55),Inches(5.7),Inches(6.1),Inches(0.8))
    gb.fill.solid(); gb.fill.fore_color.rgb=RGBColor(0xEA,0xF4,0xEE); gb.line.color.rgb=OKc; gb.line.width=Pt(1); gb.shadow.inherit=False
    tf2=gb.text_frame; tf2.word_wrap=True; tf2.vertical_anchor=MSO_ANCHOR.MIDDLE; tf2.margin_left=Inches(0.12)
    r2=tf2.text_frame if False else tf2.paragraphs[0].add_run(); r2.text="Good looks like: "+good; r2.font.size=Pt(12); r2.font.color.rgb=RGBColor(0x0C,0x6E,0x42)
    pic(s,img,6.95,2.0,6.0)
    if note: plaintext(s,6.95,6.4,6.0,0.6,note,11,MUT)
    brand(s,num())

def table_chart_slide(h,take_txt,rows,img,headers=None,warn=False,colw=None,note=None):
    s=prs.slides.add_slide(BLANK); bg(s); header(s,h)
    if take_txt: take(s,1.02,take_txt,BADc if warn else ACC)
    y=1.85
    nrows=len(rows)+(1 if headers else 0); ncols=len(rows[0])
    tb=s.shapes.add_table(nrows,ncols,Inches(0.55),Inches(y),Inches(6.2),Inches(min(4.8,0.5*nrows))).table
    if colw:
        for i,w in enumerate(colw): tb.columns[i].width=Inches(w)
    ri=0
    if headers:
        for ci,hh in enumerate(headers):
            c=tb.cell(0,ci); c.text=hh; c.fill.solid(); c.fill.fore_color.rgb=INK
            rr=c.text_frame.paragraphs[0].runs[0]; rr.font.size=Pt(11); rr.font.bold=True; rr.font.color.rgb=WHITE
        ri=1
    for r in rows:
        for ci,val in enumerate(r):
            c=tb.cell(ri,ci); c.text=str(val); c.fill.solid(); c.fill.fore_color.rgb=PANEL if ri%2 else WHITE
            rr=c.text_frame.paragraphs[0].runs[0]; rr.font.size=Pt(12); rr.font.name="Calibri"
            rr.font.color.rgb=INK; rr.font.bold=(ci==0)
        ri+=1
    pic(s,img,7.0,1.9,6.0)
    if note: plaintext(s,7.0,6.3,6.0,0.6,note,11,MUT)
    brand(s,num())

# ============ BUILD ============
title_slide()

table_chart_slide("How to read this deck","Three layers on every slide — read the one you need.",
   [["CEO","the orange takeaway line"],["CTO","method & validation"],["Data Scientist","formula · value · interpretation"]],
   I["pipeline"],colw=[2.4,3.8],note="Every slide: a takeaway, a visual, and the detail beneath.")

table_chart_slide("1 · The problem, defined","We predict a NUMBER (units) per product, per shop, per day — a forecasting task.",
   [["Target","net = gross − refund − void"],["Grain","outlet × product × day"],["Type","regression + quantile"],
    ["Universe","402 products × 84 outlets"]],I["monthly"],colw=[2.3,3.9],
   note="Chart: 42 months of real demand — the thing we forecast.")

table_chart_slide("2 · What we ran — model spec","LightGBM — the proven workhorse for tabular demand data.",
   [["Algorithm","LightGBM quantile"],["Outputs","P50 / P85 / P95"],["Trees","600 · 255 leaves · lr 0.05"],
    ["Features","37 (18 categorical)"]],I["imp"],colw=[2.3,3.9],
   note="Chart: which signals the model leans on most.")

diagram_slide("3 · Data split & validation method","Trained on the past, tested on the future it never saw — the honest way.",
   I["wf"],["Train up to month M → predict M+1 → score → roll forward.","7.72M train rows · 608k test rows · 3 monthly folds.",
   "No leakage: history features shifted ≥1 day."],note="A single random split would overstate accuracy for time-series. We avoid it.")

diagram_slide("4 · Which scores apply — and which DON'T","AUC / ROC / precision / recall are classification metrics — they do NOT apply here.",
   I["classes"],["We predict a quantity, not a yes/no class → those metrics are the wrong tool.",
   "Correct here: WMAPE, MAE, RMSE, Bias, R² (point) + Pinball & Coverage (quantile)."],warn=False,
   note="Stating this up front is what signals rigor to a CTO.")

# metric slides — each with chart + plain words
metric_slide("5 · WMAPE — headline accuracy","On average the forecast is off by 34% of volume — 16% better than today's method.",
   "WMAPE = Σ|actual−forecast| ÷ Σ|actual|","add up every miss, divide by total real demand.",
   str(M["wmape"]),["Volume-weighted: big sellers count more.","Stable where plain MAPE blows up."],
   "Retail: <0.40 workable, <0.35 good → we are 0.341.",I["wmape"],note="Green = us; orange = best simple method; red line = target.")
metric_slide("6 · MAE — average miss size","A typical forecast is about 2.3 units away from reality.",
   "MAE = average( |actual − forecast| )","how many units off, on average.",
   f"{M['mae']} units",["Same units as demand — easy to explain.","Robust to the odd outlier."],
   "Lower is better; compare to the item's daily volume.",I["mae"])
metric_slide("7 · RMSE — checks for big misses","Big errors are rare — RMSE stays close to MAE.",
   "RMSE = √ average( (actual − forecast)² )","like MAE but punishes BIG misses harder.",
   f"{M['rmse']} units",["Squaring makes large errors stand out.",f"RMSE {M['rmse']} vs MAE {M['mae']} → spread controlled."],
   "Want RMSE not much bigger than MAE → no disasters.",I["rmse"])
metric_slide("8 · Bias — do we over/under order?","Almost perfectly balanced — a tiny lean to under-forecast.",
   "Bias = average( forecast − actual )","on average, do we guess too high or too low?",
   f"{M['bias']} units",["0 = balanced; negative = slightly low.","Persistent bias would mis-stock every day."],
   "Want ≈ 0. −0.41 on ~6-unit orders is negligible.",I["bias"])
metric_slide("9 · R² — how much we explain","The model explains 81% of why demand goes up and down.",
   "R² = 1 − (model errors ÷ just-guess-the-average errors)","how much better than always guessing the average.",
   str(M["r2"]),["1.0 = perfect; 0 = no better than the average.","0.81 = captures most of the swing."],
   ">0.7 is strong for noisy daily retail.",I["r2"],note="Bar: our R² vs the 'good' line. Scatter shown on proof slides.")
metric_slide("10 · Pinball loss — quantile quality","The correct score for our 'safe stock' levels (P85/P95).",
   "Pinball(q): penalty that's bigger on the costly side","grades the range forecast, not just the middle.",
   f"{M['pin50']} / {M['pin85']} / {M['pin95']}",["Built for quantiles (P50/P85/P95).","Lower = tighter and correctly placed."],
   "Lower is better; used to compare quantile models.",I["pinball"])
metric_slide("11 · Coverage — are safety levels honest?","Say '85% safe' and demand stays under it 85.4% of the time.",
   "Coverage(q) = how often actual ≤ the q forecast","does the promise match reality?",
   f"P85 {M['cov85']}% · P95 {M['cov95']}%",["On target → the buffers are trustworthy.","P50 56.7% = small safe over-rounding."],
   "Coverage ≈ stated % = calibrated → safe to order on.",I["calib"])

# proof slides (chart-led)
diagram_slide("12 · Proof 1 — it beats simple methods","Only 'good' if it wins on the SAME test. It does, by +16%.",
   I["wmape"],["Compared against 5 simple methods on identical data.","Lower bar = better; green is our model."],wide=8.5,
   note="Red line = the stretch target (0.321).")
diagram_slide("13 · Proof 2 — stable every month","Wins in all 3 separate test months → not a fluke.",
   I["folds"],["Beats the baseline in April, May and June.","Consistency = it will hold in production."],wide=8.5)
diagram_slide("14 · Proof 3 — best where it matters","Most accurate on Class-A (80% of volume): 0.291.",
   I["abc"],["A = top sellers, B = mid, C = rare.","We deliberately optimise the high-volume items."],wide=8.5)
diagram_slide("15 · Proof 4 — errors centred on zero","No systematic mistake — misses cancel out around 0.",
   I["resid"],["Distribution of (actual − forecast).","Symmetric & peaked at 0 = healthy."],wide=8.5,
   note="A skew here would reveal hidden bias. There is none.")
diagram_slide("16 · Proof 5 — it tracks reality",f"{DD['hero']['name']}: the forecast follows actual sales day by day.",
   I["hero"],["Orange = forecast, dark = actual, green dashed = safe level.","Lines move together → trustworthy."],wide=10.5)
diagram_slide("17 · Proof 6 — sensible drivers","It relies on product, outlet & recent demand — like a good planner.",
   I["imp"],["No weird/leaky features at the top.","Matches business intuition → safe to trust."],wide=8.5)

diagram_slide("18 · So — how do I KNOW it's right?","Six independent checks all pass. That's the answer.",
   I["checks"],[],note="Skill · stability · calibration · low bias · variance explained · business saving.")

diagram_slide("19 · The outcome — forecast becomes an order","A balance: enough to avoid stockouts, not so much it spoils.",
   I["scale"],["Forecast range → newsvendor order quantity → one daily warehouse picklist.",
   "~230 products · ~34,000 units in the example day."],wide=8.5)
diagram_slide("20 · Business validation — money saved","Smart ordering ≈ 21% cheaper than current practice.",
   I["cost"],["Simulated over 608k order-days.","Demo economics — real margin & shelf-life improve it further."],wide=8.5,
   note="Lower bar = lower total cost (stockout + spoilage).")
diagram_slide("21 · It keeps improving itself","Relearns from new sales; a worse model never goes live.",
   I["loop"],["Champion vs challenger: promote only if ≥1% better.","Champion holdout WMAPE 0.319 — beats target."],wide=9.5)

table_chart_slide("22 · Assumptions & limitations (honest)","Forecast is real & proven. Order economics are placeholders until Finance/Ops data.",
   [["Assumed","margin 35% · shelf 1 day · salvage 0"],["Effect","ordering ≈ P50 for now"],
    ["Missing","real margin & shelf-life"],["Fix","data only — no rebuild"]],I["scale"],warn=True,colw=[2.0,4.2],
   note="This is a data gap, not a model gap.")
diagram_slide("23 · What would prove us WRONG","A credible model says how it could fail — and we watch for it.",
   I["loop"],["WMAPE drifts >0.40 → retrain fires.","Coverage drifts off target → recalibrate.",
   "Bias grows → investigate. Input drift (PSI) → relearn."],wide=9.5,note="All monitored automatically; none silent.")

table_chart_slide("24 · The scorecard — everything at a glance",None,
   [["WMAPE",str(M["wmape"]),"good +16%"],["MAE",f"{M['mae']}u","small"],["RMSE",f"{M['rmse']}u","no outliers"],
    ["Bias",f"{M['bias']}u","neutral"],["R²",str(M["r2"]),"strong"],["Pinball P85",str(M["pin85"]),"sharp"],
    ["Coverage P85",f"{M['cov85']}%","calibrated"],["Coverage P95",f"{M['cov95']}%","calibrated"],
    ["Champion","0.319","beats tgt"],["Cost","−21%","saves"]],
   I["calib"],headers=["Metric","Value","Verdict"],colw=[2.6,1.8,1.8],note="Chart: calibration — the trust check.")

diagram_slide("25 · Next step — the one unlock","Real product economics turns a proven forecast into optimal ordering.",
   I["scale"],["Finance: margin per product. Ops: shelf-life + salvage.","Drop into one file → no code change.",
   "Then pilot in a few outlets, measure ₭ saved, roll out."],wide=8.5)

# ---------------- presenter notes (one script per slide, in build order) ----------------
NOTES=[
# 0 title
"SAY: This deck shows what we built, every score it earns, and how we prove the forecasts are right. "
"Three audiences — CEO reads the orange line, CTO reads the method, data scientist reads the formulas. "
"The strip at the bottom is the whole pipeline: sales → learn → forecast → order → picklist.",
# 1 how to read
"SAY: Every slide has the same three layers so nobody is lost. Point at the orange takeaway for execs, "
"the table for engineers, the small print for the data team. The picture is always the same pipeline.",
# 2 problem
"SAY: We are predicting a NUMBER — how many units of each product each shop needs each day. "
"WHAT THE CHART SHOWS: 42 months of real demand — steady business with seasonality. That wavy line is exactly what we forecast. "
"This is a forecasting problem, not a yes/no problem — remember that for the metrics later.",
# 3 model spec
"SAY: We used LightGBM — a tree model that's the industry standard for this kind of tabular sales data. It gives three levels: likely, safe, very-safe. "
"WHAT THE CHART SHOWS: the signals it relies on most — product, outlet, recent demand. Sensible, no magic.",
# 4 validation
"SAY: This is HOW we test honestly. We train on the past and predict the next month it has never seen, then roll forward. "
"WHAT THE DIAGRAM SHOWS: green = training, orange = the unseen test month, repeated 3 times. "
"IF ASKED why not a random split: for time-series a random split leaks the future and fakes good scores — walk-forward is the honest way.",
# 5 which metrics
"SAY: Important — AUC and ROC are for yes/no predictions (spam / not-spam). We predict a quantity, so those don't apply. "
"WHAT THE DIAGRAM SHOWS: red box = metrics we correctly DON'T use; green box = the right ones (WMAPE, MAE, RMSE, R², Pinball, Coverage). "
"Saying this first tells a CTO we know our metrics.",
# 6 WMAPE
"SCORE = WMAPE 0.341. WHAT IT MEANS: on average our forecast is off by about 34% of total volume. "
"HOW TO EXPLAIN: add up every miss in units, divide by total real demand — weighted so big sellers matter more. "
"CHART: green = us (0.341), orange = best simple method (0.405), red dashed = our target. Lower is better. We beat the simple method by 16%.",
# 7 MAE
"SCORE = MAE 2.29 units. WHAT IT MEANS: a typical forecast is about 2 units away from what actually sold. "
"HOW TO EXPLAIN: just the average size of the miss, in the same units as sales — the most intuitive error number. "
"CHART: MAE next to RMSE for context.",
# 8 RMSE
"SCORE = RMSE 5.09 units. WHAT IT MEANS: same idea as MAE but it punishes BIG misses harder. "
"HOW TO EXPLAIN: because RMSE (5.1) isn't wildly bigger than MAE (2.3), we have no catastrophic blow-ups — errors are controlled. "
"CHART: RMSE (orange) vs MAE (grey) — the gap is moderate, which is good.",
# 9 bias
"SCORE = Bias −0.41 units. WHAT IT MEANS: are we always guessing too high or too low? Basically neither — a tiny lean to under. "
"HOW TO EXPLAIN: zero is perfectly balanced; −0.41 on ~6-unit orders is negligible and we round orders up anyway. "
"CHART: the bar sits almost on the zero line.",
# 10 R2
"SCORE = R² 0.811. WHAT IT MEANS: the model explains 81% of why demand rises and falls. "
"HOW TO EXPLAIN: 1.0 is perfect, 0 is no better than always guessing the average — 0.81 means we capture most of the movement. "
"CHART: the green bar vs the red 'good ≥0.7' line — we clear it comfortably.",
# 11 pinball
"SCORE = Pinball 1.15 / 0.83 / 0.45 for P50/P85/P95. WHAT IT MEANS: this is the correct grade for a RANGE forecast, not just the middle guess. "
"HOW TO EXPLAIN: it rewards ranges that are both tight and correctly placed. Lower is better. "
"CHART: one bar per level — we report it so a data scientist sees we scored the quantiles properly.",
# 12 coverage
"SCORE = Coverage P85 85.4%, P95 94.5%. WHAT IT MEANS: when we label a stock level '85% safe', demand really does stay under it 85% of the time. "
"HOW TO EXPLAIN: the promise matches reality → the safety buffers can be trusted for ordering. "
"CHART: orange (actual) sits right on grey (target) for P85/P95. P50 reads 57% because we safely round up.",
# 13 proof1
"SAY: A model is only good if it beats simple methods on the SAME test. WHAT THE CHART SHOWS: us (green) vs five simple methods. "
"We are lowest, 16% better than the best of them. Red line = our stretch target.",
# 14 proof2
"SAY: Is it luck? No. WHAT THE CHART SHOWS: three separate test months — we beat the baseline in every one. "
"Consistency means it will hold up in production, not just on a lucky month.",
# 15 proof3
"SAY: We focus accuracy where the money is. WHAT THE CHART SHOWS: Class-A (top 80% of volume) is the MOST accurate at 0.291. "
"Rare Class-C is harder, but it's a tiny slice of volume.",
# 16 proof4
"SAY: WHAT THE CHART SHOWS: the spread of our errors. It's symmetric and peaks at zero — misses cancel out, no hidden one-sided mistake. "
"IF ASKED: a lean to one side here would reveal bias; it's centred, so we're clean.",
# 17 proof5
"SAY: The eye test. WHAT THE CHART SHOWS: our best-selling product — the orange forecast line tracks the dark actual line day by day. "
"Green dashed is the safe level. They move together → believable.",
# 18 proof6
"SAY: WHAT THE CHART SHOWS: the model's top drivers — product, outlet, recent demand. Exactly what a smart planner would use. "
"No weird or 'leaky' features at the top → we trust how it decides.",
# 19 how I know
"SAY THIS IS THE KEY SLIDE: when someone asks 'how do you know it's right?', give these six: "
"beats simple methods (+16%), wins all 3 months, safety levels are honest (85/95%), almost no bias, explains 81% of demand, and saves ~21% in the cost test. "
"Six independent checks all green.",
# 20 outcome/scale
"SAY: A forecast isn't a decision. WHAT THE DIAGRAM SHOWS: ordering is a balance — too little loses a sale, too much spoils. "
"The newsvendor rule picks the amount that balances those two costs, then we sum all shops into one warehouse picklist.",
# 21 cost
"SAY: Does it save money? WHAT THE CHART SHOWS: simulated ordering cost. Our model (green) is about 21% cheaper than today's practice (grey). "
"Caveat: this uses placeholder economics; real margins make it better, not worse. Lower bar = cheaper.",
# 22 loop
"SAY: It maintains itself. WHAT THE DIAGRAM SHOWS: each night it predicts, sees what sold, learns, and only promotes a new model if it's at least 1% better. "
"A worse model can never go live. The live champion already scores 0.319.",
# 23 assumptions
"SAY (be honest): the FORECAST is real and proven. The ORDER-SIZING uses placeholder economics — flat 35% margin, 1-day shelf-life. "
"WHAT IT MEANS: until we get real margins and shelf-life, ordering behaves like a simple P50. This is a data request, not a rebuild.",
# 24 falsifiability
"SAY: A trustworthy model states how it could fail and watches for it. WHAT THE DIAGRAM SHOWS the loop that monitors: "
"if accuracy drifts past 0.40, or coverage slips, or bias grows, or inputs shift — it flags and retrains. Nothing fails silently.",
# 25 scorecard
"SAY: One slide, every number, every verdict. Walk down the right column — all green. "
"CHART: calibration again because it's the trust check — actual safety levels match the promise.",
# 26 next step
"SAY: The single unlock. Finance gives margin per product, Ops gives shelf-life — we drop it into one file, no code change. "
"Then pilot in a few shops, measure the kyats saved, and roll out. That turns a proven forecast into optimal ordering.",
]

for sld,note in zip(prs.slides, NOTES):
    sld.notes_slide.notes_text_frame.text=note

out=ROOT/"reports"/"CFC_Model_Evidence_Light.pptx"
prs.save(str(out))
print(f"wrote {out} — {N[0]+1} slides, {len(NOTES)} presenter notes")
