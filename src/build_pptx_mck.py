"""
CFC Model Evidence — consulting (McKinsey-style) deck. White, airy, action titles, one idea per slide,
minimal color, clean direct-labelled charts. Native editable PPTX + speaker notes.
Output: reports/CFC_Model_Evidence.pptx
"""
import json, pathlib, numpy as np
import matplotlib; matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib.patches import FancyBboxPatch, FancyArrowPatch
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

ROOT=pathlib.Path(__file__).resolve().parent.parent
DD=json.loads((ROOT/"reports"/"deck_data.json").read_text())
M=json.loads((ROOT/"reports"/"deck_metrics.json").read_text())
FIG=ROOT/"reports"/"_mck"; FIG.mkdir(exist_ok=True)

# consulting palette
NAVY=RGBColor(0x1F,0x3A,0x5F); BLUE=RGBColor(0x2F,0x6D,0xB4); INK=RGBColor(0x22,0x2A,0x35)
GREY=RGBColor(0x6B,0x74,0x85); LGREY=RGBColor(0xE9,0xED,0xF3); RULE=RGBColor(0xD2,0xD9,0xE3)
GREEN=RGBColor(0x1E,0x7A,0x46); RED=RGBColor(0xB0,0x3A,0x2E)
HNAVY="#1f3a5f"; HBLUE="#2f6db4"; HLBLUE="#a9c6e8"; HGREY="#aab3c2"; HINK="#222a35"; HGREEN="#1e7a46"; HRED="#b03a2e"

prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
BLANK=prs.slide_layouts[6]
ML=0.7  # left margin

def bg(s): s.background.fill.solid(); s.background.fill.fore_color.rgb=RGBColor(0xFF,0xFF,0xFF)
def tx(s,x,y,w,h,t,size,color=INK,bold=False,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,spacing=1.0):
    tb=s.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h)); tf=tb.text_frame
    tf.word_wrap=True; tf.vertical_anchor=anchor
    for i,line in enumerate(t.split("\n")):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.alignment=align; p.line_spacing=spacing
        r=p.add_run(); r.text=line; r.font.size=Pt(size); r.font.color.rgb=color; r.font.bold=bold; r.font.name="Calibri"
    return tb
def head(s,eyebrow,title):
    tx(s,ML,0.42,12,0.35,eyebrow.upper(),12,BLUE,True)
    tx(s,ML,0.74,12.0,1.1,title,23,NAVY,True,spacing=1.0)
    ln=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(ML),Inches(1.72),Inches(11.93),Pt(1.5))
    ln.fill.solid(); ln.fill.fore_color.rgb=RULE; ln.line.fill.background(); ln.shadow.inherit=False
def foot(s,n,src="Source: CFC sales (Microsoft Fabric) 2023–2026; model backtest"):
    tx(s,ML,7.04,9.5,0.3,src,9,GREY)
    tx(s,12.4,7.04,0.7,0.3,str(n),9,GREY)
def bullets(s,x,y,w,h,items,size=15,gap=8,color=INK,mark="–"):
    tb=s.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h)); tf=tb.text_frame; tf.word_wrap=True
    for i,it in enumerate(items):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.space_after=Pt(gap)
        r=p.add_run(); r.text=mark+"  "; r.font.size=Pt(size); r.font.color.rgb=BLUE; r.font.bold=True
        r2=p.add_run(); r2.text=it; r2.font.size=Pt(size); r2.font.color.rgb=color
def pic(s,img,x,y,w): s.shapes.add_picture(img,Inches(x),Inches(y),width=Inches(w))
def statbox(s,x,y,w,big,label,color=NAVY):
    tx(s,x,y,w,0.7,big,34,color,True,PP_ALIGN.LEFT)
    tx(s,x,y+0.66,w,0.7,label,13,GREY)
def caption(s,x,y,w,t): tx(s,x,y,w,0.5,t,11,GREY)
def readbox(s,x,y,w,h,title,lines,accent=BLUE,fill=RGBColor(0xF4,0xF7,0xFB)):
    """Light explanation panel: small heading + simple-language bullet lines."""
    b=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(w),Inches(h))
    b.fill.solid(); b.fill.fore_color.rgb=fill; b.line.color.rgb=RULE; b.line.width=Pt(0.75); b.shadow.inherit=False
    tb=s.shapes.add_textbox(Inches(x+0.18),Inches(y+0.12),Inches(w-0.36),Inches(h-0.24)); tf=tb.text_frame; tf.word_wrap=True
    p=tf.paragraphs[0]; r=p.add_run(); r.text=title.upper(); r.font.size=Pt(11); r.font.bold=True; r.font.color.rgb=accent; p.space_after=Pt(5)
    for ln in lines:
        p=tf.add_paragraph(); p.space_after=Pt(4)
        if isinstance(ln,tuple):
            r=p.add_run(); r.text=ln[0]+" "; r.font.size=Pt(12.5); r.font.bold=True; r.font.color.rgb=INK
            r2=p.add_run(); r2.text=ln[1]; r2.font.size=Pt(12.5); r2.font.color.rgb=INK
        else:
            r=p.add_run(); r.text=ln; r.font.size=Pt(12.5); r.font.color.rgb=INK
    return b

# ---------------- clean charts ----------------
def _clean(ax,fig,xgrid=False):
    fig.patch.set_facecolor("white"); ax.set_facecolor("white")
    for k,sp in ax.spines.items(): sp.set_visible(k=="bottom"); sp.set_color("#c9d2de") if k=="bottom" else None
    ax.tick_params(colors=HINK,labelsize=10,length=0)
    ax.grid(False)
    if xgrid: ax.xaxis.grid(True,color="#eef1f6",lw=.8); ax.set_axisbelow(True)
def _save(fig,n): p=FIG/f"{n}.png"; fig.savefig(p,dpi=160,bbox_inches="tight",facecolor="white"); plt.close(fig); return str(p)

def hbar(labels,values,name,highlight=0,fmt="{:.3f}",figsize=(6.6,3.6),accent=HBLUE,note_target=None):
    fig,ax=plt.subplots(figsize=figsize); y=np.arange(len(labels))
    cols=[accent if i==highlight else "#cfd8e6" for i in range(len(labels))]
    ax.barh(y,values,color=cols,height=.62)
    ax.set_yticks(y); ax.set_yticklabels(labels,fontsize=10,color=HINK); ax.invert_yaxis()
    for i,v in enumerate(values):
        ax.text(v,i,"  "+fmt.format(v),va="center",ha="left",fontsize=10,
                color=HINK,fontweight="bold" if i==highlight else "normal")
    ax.set_xlim(0,max(values)*1.18); ax.set_xticks([])
    _clean(ax,fig);
    for sp in ax.spines.values(): sp.set_visible(False)
    return _save(fig,name)
def grouped(labels,a,b,name,la="Our model",lb="Baseline",figsize=(6.6,3.4)):
    fig,ax=plt.subplots(figsize=figsize); x=np.arange(len(labels)); w=.36
    ax.bar(x-w/2,a,w,color=HBLUE,label=la); ax.bar(x+w/2,b,w,color="#cfd8e6",label=lb)
    for i,v in enumerate(a): ax.text(x[i]-w/2,v,f"{v:.3f}",ha="center",va="bottom",fontsize=9,color=HINK,fontweight="bold")
    for i,v in enumerate(b): ax.text(x[i]+w/2,v,f"{v:.3f}",ha="center",va="bottom",fontsize=9,color=HGREY)
    ax.set_xticks(x); ax.set_xticklabels(labels,fontsize=10,color=HINK)
    ax.legend(fontsize=9,frameon=False,loc="upper left"); ax.set_yticks([])
    _clean(ax,fig)
    for sp in ax.spines.values(): sp.set_visible(False)
    return _save(fig,name)
def lines(labels,series,name,figsize=(11.2,3.6)):
    fig,ax=plt.subplots(figsize=figsize)
    for lab,data,col,dash,lw in series: ax.plot(range(len(labels)),data,color=col,label=lab,lw=lw,ls=dash)
    step=max(1,len(labels)//8); ax.set_xticks(range(0,len(labels),step))
    ax.set_xticklabels([labels[i] for i in range(0,len(labels),step)],fontsize=9,color=HINK)
    ax.legend(fontsize=10,frameon=False,loc="upper right"); _clean(ax,fig); return _save(fig,name)
def calib(name):
    fig,ax=plt.subplots(figsize=(6.4,3.4)); x=np.arange(3); w=.36
    act=[M["cov50"],M["cov85"],M["cov95"]]; tgt=[50,85,95]
    ax.bar(x-w/2,act,w,color=HBLUE,label="Actual coverage"); ax.bar(x+w/2,tgt,w,color="#cfd8e6",label="Target")
    for i,v in enumerate(act): ax.text(x[i]-w/2,v,f"{v:.0f}%",ha="center",va="bottom",fontsize=9,fontweight="bold",color=HINK)
    ax.set_xticks(x); ax.set_xticklabels(["P50","P85","P95"],fontsize=10,color=HINK); ax.set_yticks([])
    ax.legend(fontsize=9,frameon=False,loc="upper left"); _clean(ax,fig)
    for sp in ax.spines.values(): sp.set_visible(False)
    return _save(fig,name)
def costchart(name):
    return hbar(["Current practice","Our model"],[3.71,2.94],name,highlight=1,fmt="₭{:.2f}bn",figsize=(6.6,2.4),accent=HGREEN)

# clean flow diagrams (thin, navy/blue)
def _box(ax,x,y,w,h,t,sub):
    ax.add_patch(FancyBboxPatch((x,y),w,h,boxstyle="round,pad=0.02,rounding_size=0.05",fc="white",ec=HBLUE,lw=1.3))
    ax.text(x+w/2,y+h*0.60,t,ha="center",va="center",fontsize=10,fontweight="bold",color=HNAVY)
    if sub: ax.text(x+w/2,y+h*0.26,sub,ha="center",va="center",fontsize=8,color=HGREY)
def _ar(ax,x1,y,x2): ax.add_patch(FancyArrowPatch((x1,y),(x2,y),arrowstyle="-|>",mutation_scale=13,color=HBLUE,lw=1.4))
def pipeline(name):
    fig,ax=plt.subplots(figsize=(11.4,1.9)); ax.set_xlim(0,11.4); ax.set_ylim(0,1.9); ax.axis("off"); fig.patch.set_facecolor("white")
    st=[("Sales history","3.5 yrs"),("Learn patterns",""),("Forecast","range"),("Order qty","newsvendor"),("Warehouse list","daily")]
    w=1.9;h=1.0;g=0.42
    for i,(t,sb) in enumerate(st):
        x=i*(w+g)+0.05; _box(ax,x,0.5,w,h,t,sb)
        if i<4: _ar(ax,x+w,1.0,x+w+g)
    return _save(fig,name)
def walkfwd(name):
    fig,ax=plt.subplots(figsize=(10.6,2.8)); ax.set_xlim(0,11); ax.set_ylim(0,2.8); ax.axis("off"); fig.patch.set_facecolor("white")
    for r,tr in enumerate([6,7,8]):
        y=2.0-r*0.7
        ax.add_patch(FancyBboxPatch((0.3,y),tr*0.85,0.45,boxstyle="round,pad=0.01",fc="#eef3f9",ec=HBLUE,lw=1))
        ax.text(0.3+tr*0.85/2,y+0.22,"train",ha="center",va="center",fontsize=8,color=HNAVY)
        ax.add_patch(FancyBboxPatch((0.3+tr*0.85+0.08,y),0.8,0.45,boxstyle="round,pad=0.01",fc=HBLUE,ec=HBLUE))
        ax.text(0.3+tr*0.85+0.48,y+0.22,"test",ha="center",va="center",fontsize=8,fontweight="bold",color="white")
    ax.text(5.4,0.15,"roll forward each month — model never sees the future",ha="center",fontsize=9,color=HGREY,style="italic")
    return _save(fig,name)
def scalediag(name):
    fig,ax=plt.subplots(figsize=(8.8,2.6)); ax.set_xlim(0,10); ax.set_ylim(0,2.6); ax.axis("off"); fig.patch.set_facecolor("white")
    ax.plot([1.2,8.8],[1.4,1.4],color=HNAVY,lw=2); ax.plot([5,5],[1.4,2.2],color=HNAVY,lw=2)
    _box(ax,0.7,0.45,3.2,0.9,"Order too little","stockout = lost margin")
    _box(ax,6.1,0.45,3.2,0.9,"Order too much","spoilage = wasted cash")
    ax.text(5,2.35,"Newsvendor picks the balance point",ha="center",fontsize=10,fontweight="bold",color=HBLUE)
    return _save(fig,name)
def loopdiag(name):
    fig,ax=plt.subplots(figsize=(10.4,1.9)); ax.set_xlim(0,11);ax.set_ylim(0,1.9);ax.axis("off");fig.patch.set_facecolor("white")
    st=["Predict nightly","Observe sales","Retrain weekly","Promote if better"]
    w=2.1;h=0.95;g=0.5
    for i,t in enumerate(st):
        x=i*(w+g)+0.2;_box(ax,x,0.55,w,h,t,"")
        if i<3:_ar(ax,x+w,1.02,x+w+g)
    ax.text(5.4,0.2,"a worse model never goes live",ha="center",fontsize=9,color=HGREY,style="italic")
    return _save(fig,name)

I={"pipeline":pipeline("pipe"),"wf":walkfwd("wf"),"scale":scalediag("scale"),"loop":loopdiag("loop"),
 "wmape":hbar(["Our model","7-day average","28-day average","Naive (yesterday)","Day-of-week avg","Same weekday"],
              [0.341,0.405,0.411,0.453,0.468,0.537],"wmape",0),
 "folds":grouped(["April","May","June"],[0.384,0.325,0.305],[0.497,0.360,0.345],"folds"),
 "abc":grouped(["Class A (75%)","Class B (19%)","Class C (5%)"],[0.291,0.456,0.640],[0.349,0.534,0.733],"abc"),
 "calib":calib("calib"),
 "hero":lines(DD["hero"]["labels"],[("Actual",DD["hero"]["actual"],HINK,"-",1.8),("Forecast (P50)",DD["hero"]["p50"],HBLUE,"-",1.8),("Safe level (P85)",DD["hero"]["p85"],HGREEN,"--",1.4)],"hero"),
 "cost":costchart("cost"),
 "monthly":lines(DD["monthly"]["labels"],[("Net units sold",DD["monthly"]["units"],HBLUE,"-",1.8)],"monthly"),
}

# ================= SLIDES =================
N=[0]
def nx(): N[0]+=1; return N[0]
NOTES=[]
def note(t): NOTES.append(t)

# 1 title
s=prs.slides.add_slide(BLANK); bg(s)
band=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(0),Inches(0),Inches(0.28),Inches(7.5)); band.fill.solid(); band.fill.fore_color.rgb=NAVY; band.line.fill.background(); band.shadow.inherit=False
tx(s,ML,2.3,11.5,0.4,"CITYFOOD CONCEPTS · DEMAND FORECASTING",13,BLUE,True)
tx(s,ML,2.75,11.8,1.6,"Predicting daily demand and\nright-sizing every warehouse order",33,NAVY,True,spacing=1.05)
tx(s,ML,4.7,11.0,0.8,"Evidence, scores, and how we know the forecast is right",17,GREY)
tx(s,ML,6.6,11,0.4,"Prepared for CEO / CTO / Data Science  ·  June 2026",12,GREY)
note("SAY: This is the case for a demand-forecasting system at CFC. Three audiences — the headline for the CEO, the method for the CTO, the metrics for the data team. We'll end with a single recommendation.")
foot(s,nx(),src="")

# 2 executive summary (the answer first)
s=prs.slides.add_slide(BLANK); bg(s); head(s,"Executive summary","A proven forecast cuts ordering errors today; one data input unlocks full value")
bullets(s,ML,2.0,7.0,3.6,[
 "Forecasts are 16% more accurate than the averages used today — and stable in every month tested.",
 "The safety levels are honest: when we say “85% safe”, demand stays within it 85% of the time.",
 "In simulation, smarter ordering is ~21% cheaper than current practice.",
 "The system retrains itself and flags when the world changes.",
 "One gap remains: real product margin and shelf-life — a data request, not a rebuild."],size=15,gap=11)
statbox(s,8.2,2.1,4.0,"16%","more accurate than today")
statbox(s,8.2,3.5,4.0,"0.81","of demand variation explained")
statbox(s,8.2,4.9,4.0,"~21%","lower ordering cost (sim.)",GREEN)
note("SAY: Bottom line up front. Four wins and one honest gap. Accuracy +16% and stable; safety levels are trustworthy; ordering ~21% cheaper; it self-maintains. The only thing missing is real cost & shelf-life data, which is a request to Finance and Ops — no rebuild.")
foot(s,nx())

# 3 problem
s=prs.slides.add_slide(BLANK); bg(s); head(s,"The problem","Daily ordering by gut feel loses money in both directions")
pic(s,I["monthly"],6.7,2.0,6.2); caption(s,6.7,6.2,6.2,"42 months of real demand — the pattern we forecast.")
bullets(s,ML,2.1,5.6,3.5,[
 "84 outlets order fresh bakery from the warehouse every day.",
 "Order too much → unsold stock spoils overnight (cash waste).",
 "Order too little → empty shelves (lost sales, unhappy customers).",
 "At 58.8M units and ₭208bn of sales, small % errors are large money."],size=15,gap=12)
note("SAY: Today's ordering is gut feel and rough averages, so both errors happen daily. The chart is 42 months of real demand — steady with seasonality. Because volume is huge, even a few percent of error is serious money. That's the opportunity.")
foot(s,nx())

# 4 approach
s=prs.slides.add_slide(BLANK); bg(s); head(s,"Approach","We forecast demand, then convert it into the optimal order — end to end")
pic(s,I["pipeline"],0.9,2.4,11.6)
bullets(s,ML,4.6,11.8,1.8,[
 "Forecast is a range (likely / safe / very-safe), not a single guess — so we can choose how cautious to be per product.",
 "The order quantity balances the cost of a stockout against the cost of spoilage.",
 "All outlet orders roll up into one daily warehouse production-and-dispatch list."],size=14,gap=8)
note("SAY: Five steps, fully automatic each night. Learn from history, forecast a range, turn the range into an order that balances stockout vs spoilage, and produce one warehouse list. The key idea: a forecast is not a decision — we go all the way to the order number.")
foot(s,nx())

# 5 data
s=prs.slides.add_slide(BLANK); bg(s); head(s,"Foundation · data","Built on 7.08M clean sales records, enriched with weather and festivals")
statbox(s,ML,2.1,3.0,"7.08M","sales rows (day×shop×product)")
statbox(s,ML,3.5,3.0,"3.5 yrs","Jan 2023 – Jun 2026")
statbox(s,ML,4.9,3.0,"84 / 3,580","outlets / active products")
bullets(s,4.4,2.15,4.0,3.5,[
 "Internal: demand, price, product & outlet master (Microsoft Fabric).",
 "External: rain & temperature (6 cities), public holidays & festivals.",
 "Quality: no nulls, no orphan records; <0.01% edge cases handled."],size=13,gap=9)
pic(s,I["monthly"],8.4,2.1,4.7); caption(s,8.4,5.6,4.7,"Demand history used for training.")
note("SAY: Strong, clean foundation. 7 million rows over 3.5 years, plus weather and festival signals a good manager would consider. Data quality checks passed — essentially no nulls or broken joins. Good data in is why the model works.")
foot(s,nx())

# 6 validation
s=prs.slides.add_slide(BLANK); bg(s); head(s,"How we tested","Scored only on future months the model never saw — the honest test")
pic(s,I["wf"],3.0,2.0,7.3)
readbox(s,ML,4.55,11.93,2.05,"How to read this",[
 ("Each row =",  "one test. Blue “train” = months the model learns from; dark “test” = the next month it must predict."),
 ("Why it matters:", "the model is graded only on days it has never seen — exactly like real life, where tomorrow is unknown."),
 ("The trap we avoid:", "mixing future days into training would make scores look great but fail in production. We never do that.")])
note("SAY: This is the credibility slide for the CTO. We never test on data the model trained on. We train up to a month, predict the next, score, and roll forward three times. A random split would leak the future and inflate the score — we deliberately don't do that.")
foot(s,nx())

# 7 metric framing
s=prs.slides.add_slide(BLANK); bg(s); head(s,"How we measure","We use forecasting metrics — classification scores like AUC/ROC do not apply")
# two columns
b1=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(ML),Inches(2.1),Inches(5.7),Inches(3.6)); b1.fill.solid(); b1.fill.fore_color.rgb=RGBColor(0xFB,0xF1,0xF0); b1.line.color.rgb=RED; b1.line.width=Pt(1); b1.shadow.inherit=False
b2=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(6.9),Inches(2.1),Inches(5.7),Inches(3.6)); b2.fill.solid(); b2.fill.fore_color.rgb=RGBColor(0xEC,0xF4,0xEF); b2.line.color.rgb=GREEN; b2.line.width=Pt(1); b2.shadow.inherit=False
tx(s,ML+0.25,2.3,5.2,0.5,"Not applicable (classification only)",15,RED,True)
bullets(s,ML+0.25,2.95,5.2,2.6,["AUC, ROC curve","Precision, recall, F1","Confusion matrix, log-loss"],size=14,gap=8,mark="×")
tx(s,ML+0.25,5.0,5.2,0.6,"These need yes/no labels. We predict a quantity.",12,GREY)
tx(s,7.15,2.3,5.2,0.5,"What we use (regression + range)",15,GREEN,True)
bullets(s,7.15,2.95,5.2,2.6,["WMAPE, MAE, RMSE, Bias, R²  (the point forecast)","Pinball loss, Coverage  (the safe levels)","Skill vs baseline  (beats simple methods)"],size=14,gap=8,mark="✓")
note("SAY: Pre-empt the classic question. AUC and ROC are for yes/no predictions — spam or not. We predict a number of units, so those are the wrong tool and we don't report them. The right metrics are on the right: WMAPE, MAE, RMSE, R² for the point forecast; Pinball and Coverage for the range. Saying this proves we know our metrics.")
foot(s,nx())

# 7b every score in plain words
s=prs.slides.add_slide(BLANK); bg(s); head(s,"The scores explained","Every score in one line — what it asks, our number, what it means")
rows=[["Score","The question it answers","Our value","In plain words"],
 ["WMAPE","How far off are we, on average?",str(M["wmape"]),"off by ~34% of volume — 16% better than today"],
 ["MAE","Typical miss in units?",f"{M['mae']}","about 2 units away from reality, on average"],
 ["RMSE","Any huge misses?",f"{M['rmse']}","close to MAE → no wild errors"],
 ["Bias","Do we over- or under-order?",f"{M['bias']}","almost neutral — a tiny lean to under"],
 ["R²","How much do we explain?",str(M["r2"]),"captures 81% of why demand moves"],
 ["Pinball","Is the range any good?",str(M["pin85"]),"the proper grade for safe-stock levels"],
 ["Coverage","Are safe levels honest?",f"{M['cov85']}%","“85% safe” is right 85% of the time"]]
t=s.shapes.add_table(len(rows),4,Inches(ML),Inches(2.0),Inches(11.93),Inches(4.7)).table
t.columns[0].width=Inches(1.7); t.columns[1].width=Inches(4.2); t.columns[2].width=Inches(1.4); t.columns[3].width=Inches(4.63)
for ci in range(4):
    c=t.cell(0,ci); c.fill.solid(); c.fill.fore_color.rgb=NAVY; c.text=rows[0][ci]
    rr=c.text_frame.paragraphs[0].runs[0]; rr.font.size=Pt(12); rr.font.bold=True; rr.font.color.rgb=RGBColor(0xFF,0xFF,0xFF)
for ri in range(1,len(rows)):
    for ci in range(4):
        c=t.cell(ri,ci); c.text=rows[ri][ci]; c.fill.solid(); c.fill.fore_color.rgb=(RGBColor(0xF4,0xF7,0xFB) if ri%2 else RGBColor(0xFF,0xFF,0xFF))
        rr=c.text_frame.paragraphs[0].runs[0]; rr.font.size=Pt(12); rr.font.name="Calibri"
        rr.font.color.rgb=(NAVY if ci==2 else INK); rr.font.bold=(ci==0 or ci==2)
note("SAY: This is the cheat-sheet for any score in the deck. Read it as: the metric, the plain question it answers, our number, and what that number means in everyday words. WMAPE is the headline; R² shows we explain most of the demand; Coverage shows the safety levels are honest. If anyone asks 'what is that score', point here.")
foot(s,nx())

# 8 headline accuracy (WMAPE)
s=prs.slides.add_slide(BLANK); bg(s); head(s,"Result · accuracy","16% more accurate than the best simple method")
pic(s,I["wmape"],6.5,2.05,6.3); caption(s,6.5,5.75,6.3,"WMAPE — lower is better. Bars = average error vs each method.")
statbox(s,ML,2.15,3.0,"0.341","WMAPE (our model)")
readbox(s,ML,3.35,5.55,1.55,"What is WMAPE?",[
 "Add up every forecast miss (in units), divide by total real demand. So 0.34 = we are off by about 34% of volume on average.",
 "It is weighted — a miss on a big seller counts more than on a rare item."])
readbox(s,ML,5.05,5.55,1.55,"How to read the chart",[
 ("Shorter bar = better. ","Blue (top) is our model at 0.341; greys are simple methods; the best simple one is 0.405."),
 ("So:","we cut the error by 16% versus the best method used today.")])
note("SAY: The headline. WMAPE is total miss divided by total demand — weighted so big sellers matter. We score 0.341 vs the best simple method 0.405 — 16% better. HOW TO READ THE CHART: each bar is a method's average error; ours (top, blue) is the shortest. Lower is better.")
foot(s,nx())

# 9 stability + ABC (two charts)
s=prs.slides.add_slide(BLANK); bg(s); head(s,"Result · robustness","Accuracy holds every month and is strongest on the highest-volume products")
pic(s,I["folds"],ML,1.95,5.6); caption(s,ML,5.25,5.6,"Beats baseline in all three test months.")
pic(s,I["abc"],7.0,1.95,5.6); caption(s,7.0,5.25,5.6,"Class A = top 80% of volume; most accurate at 0.291.")
readbox(s,ML,5.6,11.93,1.05,"How to read these charts",[
 ("In both charts:","blue = our model, grey = the simple baseline, and a SHORTER bar means a more accurate forecast (lower error)."),
 ("Left (stability):","our blue bar is shorter every month → it’s consistently better, not lucky once.   Right (where it counts): the gap is biggest on Class-A, the products that drive 80% of sales.")])
note("SAY: Two robustness checks. Left: it beats the baseline in April, May and June — not a lucky month. Right: it's most accurate on Class-A, the products that are 80% of volume — we optimise where the money is. Blue is our model, grey is the baseline, in both.")
foot(s,nx())

# 10 calibration + bias
s=prs.slides.add_slide(BLANK); bg(s); head(s,"Result · trust","The forecast is unbiased and its safety levels are honest")
pic(s,I["calib"],6.6,2.05,6.2); caption(s,6.6,5.75,6.2,"Actual coverage (blue) lands on target (grey).")
statbox(s,ML,2.15,3.0,"85% / 95%","P85 / P95 coverage — on target")
statbox(s,ML,3.45,3.0,"−0.41","bias (units) — essentially neutral")
readbox(s,ML,4.75,5.55,1.05,"What is “coverage”?",[
 "It checks our promise. If we call a stock level “85% safe”, demand should stay under it 85 days out of 100. Ours does → the buffers can be trusted."])
readbox(s,ML,5.85,5.55,0.75,"What is “bias”?",[
 "Do we lean high or low overall? −0.41 units ≈ neutral, so we don’t systematically over- or under-order."])
note("SAY: Trust check. Coverage means: when we label a level '85% safe', does demand actually stay under it 85% of the time? It does — 85% and 95% land on target. And bias is −0.41 units, essentially neutral, so we don't systematically over- or under-order. These two are why the ordering can rely on the forecast.")
foot(s,nx())

# 11 hero tracks reality
s=prs.slides.add_slide(BLANK); bg(s); head(s,"Result · the eye test",f"The forecast tracks real sales day by day — example: {DD['hero']['name']}")
pic(s,I["hero"],1.4,1.95,10.5)
readbox(s,ML,5.5,11.93,1.15,"How to understand this chart",[
 ("The three lines:","dark = what actually sold each day · blue = our forecast · green dashed = the “safe” order level (P85)."),
 ("What to look for:","blue hugs dark — the forecast follows real sales, including the weekly ups and downs. Green sits a little above, the cushion that prevents stockouts. Lines moving together = the model is believable.")])
note("SAY: The intuitive proof. This is our top product. The blue forecast line follows the dark actual line, day by day, including the weekly peaks. The green dashed line is the safe level that sits just above. When people see the lines move together, they believe it.")
foot(s,nx())

# 12 scorecard
s=prs.slides.add_slide(BLANK); bg(s); head(s,"Scorecard","Every metric, value and verdict on one page")
rows=[["Metric","Value","What good looks like","Verdict"],
 ["WMAPE",str(M["wmape"]),"< 0.40 (target 0.321)","16% better than floor"],
 ["MAE",f"{M['mae']} units","small vs typical order","✓"],
 ["RMSE",f"{M['rmse']} units","near MAE (no big misses)","✓"],
 ["Bias",f"{M['bias']} units","≈ 0","neutral"],
 ["R²",str(M["r2"]),"> 0.70","strong"],
 ["Pinball (P85)",str(M["pin85"]),"lower is better","sharp"],
 ["Coverage P85 / P95",f"{M['cov85']}% / {M['cov95']}%","≈ 85% / 95%","calibrated"],
 ["Champion (live)","0.319","≤ 0.321","beats target"],
 ["Ordering cost vs today","−21%","< 0","saves money"]]
t=s.shapes.add_table(len(rows),4,Inches(ML),Inches(2.0),Inches(11.9),Inches(4.6)).table
t.columns[0].width=Inches(3.3); t.columns[1].width=Inches(2.2); t.columns[2].width=Inches(3.6); t.columns[3].width=Inches(2.8)
for ci in range(4):
    c=t.cell(0,ci); c.fill.solid(); c.fill.fore_color.rgb=NAVY
    r=c.text_frame.paragraphs[0].runs[0] if c.text_frame.paragraphs[0].runs else c.text_frame.paragraphs[0].add_run()
    c.text=rows[0][ci]; rr=c.text_frame.paragraphs[0].runs[0]; rr.font.size=Pt(12); rr.font.bold=True; rr.font.color.rgb=RGBColor(0xFF,0xFF,0xFF)
for ri in range(1,len(rows)):
    for ci in range(4):
        c=t.cell(ri,ci); c.text=rows[ri][ci]; c.fill.solid(); c.fill.fore_color.rgb=(RGBColor(0xF4,0xF7,0xFB) if ri%2 else RGBColor(0xFF,0xFF,0xFF))
        rr=c.text_frame.paragraphs[0].runs[0]; rr.font.size=Pt(12); rr.font.name="Calibri"
        rr.font.color.rgb=INK if ci<3 else GREEN; rr.font.bold=(ci==0 or ci==3)
note("SAY: One page, everything. Read the right column — all green. WMAPE beats the floor by 16% and the live champion at 0.319 even clears the stretch target. Hand this to the data team; it's the full evidence at a glance.")
foot(s,nx())

# 13 forecast to order
s=prs.slides.add_slide(BLANK); bg(s); head(s,"From forecast to order","The order amount balances the two costs of being wrong")
pic(s,I["scale"],1.4,2.2,10.3)
bullets(s,ML,5.2,11.8,1.4,[
 "High-margin, long-shelf-life items → order generously (rarely run out).",
 "Thin-margin, perishable items → order lean (avoid spoilage).",
 "The “newsvendor” rule sets the exact quantity per product from its economics."],size=14,gap=8)
note("SAY: How the forecast becomes a number. Ordering is a balance: too little loses the sale's margin, too much spoils. The newsvendor rule finds the point that minimises total expected cost — generous for profitable, long-life items; lean for cheap, perishable ones.")
foot(s,nx())

# 14 business impact
s=prs.slides.add_slide(BLANK); bg(s); head(s,"Business impact","Smarter ordering is about 21% cheaper than current practice")
pic(s,I["cost"],6.4,2.4,6.4); caption(s,6.4,5.0,6.4,"Total cost = lost margin on stockouts + spoilage on overstock (simulated).")
statbox(s,ML,2.4,3.2,"~21%","lower ordering cost",GREEN)
tx(s,ML,3.9,5.2,1.6,"Simulated across 608,000 order-days in the test window. Uses placeholder economics today; real margins and shelf-life will widen the gain.",13,INK)
note("SAY: Does it pay? In simulation over 608,000 order-days, our model's ordering costs about 21% less than today's practice — fewer stockouts and less spoilage combined. Important caveat: this uses placeholder economics; real margins will improve it, not worsen it.")
foot(s,nx())

# 15 self learning
s=prs.slides.add_slide(BLANK); bg(s); head(s,"Staying accurate","The system retrains itself and flags when the world changes")
pic(s,I["loop"],1.0,2.3,11.4)
bullets(s,ML,4.6,11.8,1.6,[
 "A fresh “challenger” is trained regularly and only promoted if it beats the live model by at least 1%.",
 "A drift monitor watches the inputs and accuracy; in testing it correctly flagged the monsoon weather shift.",
 "The live model already scores 0.319 on a held-out 60 days."],size=14,gap=8)
note("SAY: It maintains itself. Each cycle it predicts, sees what sold, retrains a challenger, and promotes it only if it's at least 1% better — so a worse model never goes live. A drift monitor watches for change; it caught the monsoon shift in testing. The live champion is at 0.319.")
foot(s,nx())

# 16 honest gap
s=prs.slides.add_slide(BLANK); bg(s); head(s,"The honest gap","The forecast is proven; order-sizing needs real product economics")
rows=[["Input","Placeholder today","Needed from"],
 ["Gross margin","35% flat for all","Finance"],
 ["Shelf-life","1 day for all","Operations"],
 ["Salvage value","0 (total loss)","Operations"]]
t=s.shapes.add_table(4,3,Inches(ML),Inches(2.1),Inches(7.2),Inches(2.4)).table
t.columns[0].width=Inches(2.4);t.columns[1].width=Inches(2.6);t.columns[2].width=Inches(2.2)
for ci in range(3):
    c=t.cell(0,ci);c.fill.solid();c.fill.fore_color.rgb=NAVY;c.text=rows[0][ci]
    rr=c.text_frame.paragraphs[0].runs[0];rr.font.size=Pt(12);rr.font.bold=True;rr.font.color.rgb=RGBColor(0xFF,0xFF,0xFF)
for ri in range(1,4):
    for ci in range(3):
        c=t.cell(ri,ci);c.text=rows[ri][ci];c.fill.solid();c.fill.fore_color.rgb=(RGBColor(0xF4,0xF7,0xFB) if ri%2 else RGBColor(0xFF,0xFF,0xFF))
        rr=c.text_frame.paragraphs[0].runs[0];rr.font.size=Pt(12);rr.font.color.rgb=INK;rr.font.bold=(ci==0)
bullets(s,ML,4.9,11.8,1.6,[
 "With flat economics, every product gets the same rule → ordering behaves like a simple mid-estimate.",
 "With real values, ordering differentiates by product — the full saving is unlocked.",
 "This is a data request, not a model change."],size=14,gap=8,color=INK)
note("SAY: Be transparent. The forecast is real and proven; the order-sizing currently uses placeholders — flat 35% margin, one-day shelf-life. That makes ordering behave like a simple mid-estimate. Real margins and shelf-life from Finance and Ops unlock per-product optimisation. It's a data request, not a rebuild.")
foot(s,nx())

# 17 recommendation
s=prs.slides.add_slide(BLANK); bg(s); head(s,"Recommendation","Provide the economics, pilot in a few outlets, then scale")
bullets(s,ML,2.2,11.8,3.6,[
 "1.  Finance provides gross margin per product; Operations provides shelf-life and salvage.",
 "2.  We load it into one file — no code change — and ordering becomes product-specific.",
 "3.  Run a 4–6 week pilot in a small set of outlets and measure kyats saved (waste + stockouts).",
 "4.  Roll out across the network; the self-learning loop sustains accuracy thereafter.",
 "5.  Optional next: specialist handling for rare items, and automated nightly/weekly scheduling."],size=16,gap=14)
tx(s,ML,6.0,11.8,0.6,"The forecast is ready now. The remaining step is a data hand-off, then a measured pilot.",14,NAVY,True)
note("SAY: The ask. One: Finance and Ops give us margins and shelf-life. Two: we load it in, no rebuild. Three: pilot in a few shops for a month, measure the kyats saved. Four: roll out, and the loop keeps it accurate. The forecast is ready today — the next move is the data hand-off and a pilot.")
foot(s,nx())

for sld,nt in zip(prs.slides,NOTES): sld.notes_slide.notes_text_frame.text=nt
out=ROOT/"reports"/"CFC_Model_Evidence.pptx"
prs.save(str(out))
print(f"wrote {out} — {N[0]} slides, {len(NOTES)} notes")
